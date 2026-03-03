"""
PPTX 문서 파서 (RFP가 파워포인트로 된 경우)

python-pptx로 슬라이드별 텍스트(셰이프·테이블 셀)와 테이블 목록을 추출합니다.
RFP 분석용으로 다른 파서와 동일한 포맷(raw_text, tables, sections, metadata)을 반환합니다.
"""

from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from .base_parser import BaseParser
from ..utils.logger import get_logger

logger = get_logger("pptx_parser")


class PPTXParser(BaseParser):
    """
    PPTX 문서 전용 파서. 슬라이드 텍스트 및 테이블 추출.

    지원 확장자: .pptx
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def parse(self, file_path: Path) -> Dict[str, Any]:
        """
        PPTX를 파싱하여 raw_text, tables, sections, metadata 반환.

        Args:
            file_path: PPTX 파일 경로

        Returns:
            RFP 분석용 통일 포맷 딕셔너리
        """
        logger.info(f"PPTX 파싱 시작: {file_path}")

        raw_text = self.extract_text(file_path)
        tables = self.extract_tables(file_path)
        sections = self._extract_sections_from_text(raw_text)

        result = {
            "raw_text": raw_text,
            "tables": tables,
            "sections": sections,
            "metadata": {"source": "pptx", "path": str(file_path)},
        }

        logger.info(
            f"PPTX 파싱 완료: {len(raw_text)} 문자, {len(tables)} 테이블"
        )

        return result

    def extract_text(self, file_path: Path) -> str:
        """슬라이드별로 모든 셰이프의 텍스트 추출"""
        parts = []
        try:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                row_texts.append(cell.text.strip())
                            if any(row_texts):
                                slide_texts.append(" | ".join(row_texts))
                if slide_texts:
                    parts.append(f"--- 슬라이드 {slide_idx + 1} ---\n" + "\n".join(slide_texts))
        except Exception as e:
            logger.error(f"PPTX 텍스트 추출 실패: {e}")
        return "\n\n".join(parts) if parts else ""

    def extract_tables(self, file_path: Path) -> List[Dict[str, Any]]:
        """슬라이드 내 테이블 셰이프만 추출하여 테이블 목록 반환"""
        tables = []
        try:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not shape.has_table:
                        continue
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows.append(row_data)
                    if not rows:
                        continue
                    headers = rows[0]
                    data_rows = rows[1:] if len(rows) > 1 else []
                    tables.append({
                        "table_index": len(tables),
                        "headers": headers,
                        "rows": data_rows,
                        "raw_data": rows,
                        "slide": slide_idx + 1,
                    })
        except Exception as e:
            logger.error(f"PPTX 테이블 추출 실패: {e}")
        return tables

    def _extract_sections_from_text(self, text: str) -> List[Dict[str, Any]]:
        """'--- 슬라이드 N ---' 구분자 기준으로 블록을 나누어 슬라이드 단위 섹션으로 구성."""
        sections = []
        if not text.strip():
            return sections
        blocks = text.split("--- 슬라이드 ")
        for block in blocks:
            block = block.strip()
            if not block:
                continue
            first_line = block.split("\n")[0] if "\n" in block else block
            rest = "\n".join(block.split("\n")[1:]).strip() if "\n" in block else ""
            sections.append({
                "title": first_line.rstrip(" ---"),
                "content": [line for line in rest.splitlines() if line.strip()],
                "level": 1,
            })
        return sections
