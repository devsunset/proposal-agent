"""
PPTX 슬라이드 생성기 (v3.1 - Impact-8 Framework + Win Theme)

[회사명] 레이어: Modern 스타일 PPTX 생성
python-pptx를 사용하여 슬라이드 생성

v3.1 추가:
- Executive Summary 슬라이드
- Next Step 슬라이드
- Differentiation 슬라이드
- Win Theme 배지 지원
"""

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .template_manager import TemplateManager
from ..schemas.proposal_schema import (
    BulletPoint,
    ContentExample,
    ChannelStrategy,
    CampaignPlan,
    KPIItem,
    ComparisonData,
    TableData,
)
from ..utils.logger import get_logger

logger = get_logger("pptx_generator")

# =============================================================================
# guide_template.pptx 실측 좌표 기준 상수 (10.00" × 5.625")
# Title and Content_02 레이아웃 실측값으로 보정 (2026-03-06)
# =============================================================================
SLIDE_WIDTH_INCHES = 10.00
SLIDE_HEIGHT_INCHES = 5.625
MARGIN_H = 0.357        # 실측: PH left 기준 (0.357")
MARGIN_TOP = 0.214      # 실측: PH13 상단 라벨 top
CONTENT_WIDTH = 9.289   # 실측: PH 너비 기준

# 제목 영역 (PH0 실측값)
TITLE_BOX_TOP = 0.343
TITLE_BOX_HEIGHT = 0.303

# 부제목 바 영역 (PH14 실측값)
SUBTITLE_BOX_TOP = 0.815
SUBTITLE_BOX_HEIGHT = 0.225

# 장식 선 (레이아웃 Straight Connector 실측값: T=0.702")
DECO_LINE_TOP = 0.702
DECO_LINE_HEIGHT = 0.02

# 본문 영역 (PH16 실측값: T=1.333", H=3.779")
BODY_BOX_TOP = 1.333
KEY_MSG_TOP = 4.60      # 본문 하단부 핵심메시지 바 위치 (PH16 범위 내)
KEY_MSG_HEIGHT = 0.42   # 핵심메시지 바 높이
BODY_BOX_HEIGHT = KEY_MSG_TOP - BODY_BOX_TOP - 0.10  # 본문-키메시지 간격

# 슬라이드 번호 영역 (PH12 실측값: y=5.362, h=0.190)
SLIDE_NUM_TOP = 5.362
SLIDE_NUM_LEFT = 7.396

MAX_TITLE_CHARS = 80
MAX_BULLET_CHARS = 280
MAX_KEY_MSG_CHARS = 95
MAX_BULLETS_PER_SLIDE = 10
MAX_TABLE_CELL_CHARS = 80


def _truncate(text: str, max_chars: int, suffix: str = "…") -> str:
    """텍스트가 max_chars 초과 시 자르고 suffix 붙임."""
    if not text or max_chars <= 0:
        return text or ""
    s = str(text).strip()
    if len(s) <= max_chars:
        return s
    return s[: max_chars - len(suffix)].rstrip() + suffix


# guide_template.pptx 기준 색상 상수
STYLE_COLORS = {
    "primary": RGBColor(0x14, 0x30, 0x7F),    # #14307F 메인 네이비 (accent1)
    "secondary": RGBColor(0x4C, 0x66, 0x86),  # #4C6686 미드 블루그레이 (accent2)
    "accent": RGBColor(0xBD, 0x24, 0x39),     # #BD2439 레드 강조 (accent4)
    "teal": RGBColor(0x00, 0x86, 0x82),       # #008682 청록색 (accent3)
    "dark_bg": RGBColor(0x1B, 0x2C, 0x59),   # #1B2C59 다크 네이비 (accent5)
    "dark_blue": RGBColor(0x14, 0x30, 0x7F),  # 별칭
    "sky_blue": RGBColor(0x9E, 0xB2, 0xCA),   # #9EB2CA 라이트 블루 (accent6)
    "section_dark": RGBColor(0x1B, 0x2C, 0x5A),  # #1B2C5A 섹션 다크 패널
    "white": RGBColor(255, 255, 255),
    "light": RGBColor(0xDA, 0xDA, 0xDA),      # #DADADA 라이트 그레이
    "text_dark": RGBColor(0x22, 0x22, 0x22),  # #222222
    "text_light": RGBColor(0x4C, 0x66, 0x86), # #4C6686
    "text_gray": RGBColor(0x66, 0x66, 0x66),  # #666666
}


class PPTXGenerator:
    """PPTX 슬라이드 생성기"""

    def __init__(self, template_manager: TemplateManager):
        self.template_manager = template_manager
        self.design = template_manager.design_system
        self.prs: Optional[Presentation] = None
        # 템플릿 미사용 시 True → 플레이스홀더 제거 후 좌표 기반으로만 그려 겹침/디자인 깨짐 방지
        self._no_template_mode: bool = False

    def create_presentation(self, template_name: str = "base_template") -> Presentation:
        """새 프레젠테이션 생성"""
        self.prs = self.template_manager.load_template(template_name)
        self._no_template_mode = not bool((template_name or "").strip())
        return self.prs

    def get_slide_layout(self, layout_name: str, fallback_index: int = 0):
        """안전한 레이아웃 반환 (chart/diagram 생성기에서 사용)."""
        if not self.prs or not self.prs.slide_layouts:
            raise ValueError("프레젠테이션이 초기화되지 않았습니다. create_presentation()을 먼저 호출하세요.")
        layout_idx = self.template_manager.get_layout_index(layout_name)
        n = len(self.prs.slide_layouts)
        safe_idx = min(max(0, layout_idx), n - 1) if n else fallback_index
        return self.prs.slide_layouts[safe_idx]

    def _add_slide(self, layout):
        """
        슬라이드 추가 후 레이아웃의 플레이스홀더를 자동 제거.
        '제목을 추가하려면 클릭하십시오', '부제목을 입력하십시오' 등
        기본 안내 텍스트가 생성된 PPTX에 노출되지 않도록 처리.
        """
        slide = self.prs.slides.add_slide(layout)
        try:
            for ph in list(slide.placeholders):
                elem = ph._element
                parent = elem.getparent()
                if parent is not None:
                    parent.remove(elem)
        except Exception:
            pass
        return slide

    def _clear_placeholder_text(self, ph) -> None:
        """플레이스홀더 내 기본 안내 문구('부제목을 입력해 주세요' 등) 완전 제거."""
        try:
            if not getattr(ph, "has_text_frame", False):
                return
            tf = ph.text_frame
            tf.clear()
            # clear() 후 모든 단락·런의 텍스트를 빈 문자열로 덮어써 기본 문구 미노출
            for para in list(tf.paragraphs):
                try:
                    para.text = ""
                    for run in list(para.runs):
                        run.text = ""
                except Exception:
                    pass
            # 단락이 없으면 빈 단락 하나 넣어서 placeholder가 비어 보이게
            if len(tf.paragraphs) == 0:
                p = tf.add_paragraph()
                if p is not None:
                    p.text = ""
        except Exception:
            pass

    def _add_slide_with_placeholders(self, layout):
        """
        슬라이드 추가. 템플릿 미사용 시 플레이스홀더를 제거해 좌표 기반만 사용(겹침 방지).
        템플릿 사용 시 플레이스홀더 유지, 기본 안내 텍스트만 지움.
        """
        if getattr(self, "_no_template_mode", False):
            return self._add_slide(layout)
        slide = self.prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            self._clear_placeholder_text(ph)
        return slide

    def _add_slide_for_content(self, layout_name: str = "Title and Content_02") -> tuple:
        """
        콘텐츠형 슬라이드 공통 생성: 플레이스홀더 유지 + (slide, ph_title) 반환.
        table/two-column/three-column 등에서 재사용.
        ph_title 이 None이면 _add_title_textbox() fallback을 써야 함.
        """
        layout = self._get_layout(layout_name, fallback_idx=0)
        slide = self._add_slide_with_placeholders(layout)
        ph_title = self._get_ph(slide, 0)
        return slide, ph_title

    def _set_slide_title(self, slide, ph_title, title: str) -> None:
        """PH0에 제목 채우기. ph_title이 None이면 textbox fallback."""
        if ph_title:
            self._fill_ph_text(
                ph_title, _truncate(title, MAX_TITLE_CHARS),
                font_size_pt=18, bold=True,
                color=STYLE_COLORS["primary"],
                font_name=self.template_manager.get_font_name("title"),
            )
        else:
            self._add_title_textbox(slide, _truncate(title, MAX_TITLE_CHARS))

    def _get_ph(self, slide, idx: int):
        """슬라이드에서 placeholder_format.idx 로 플레이스홀더 반환, 없으면 None."""
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.idx == idx:
                    return ph
            except Exception:
                pass
        return None

    def _fill_ph_text(self, ph, text: str, font_size_pt: float, bold: bool = False,
                      color: Optional[RGBColor] = None, font_name: Optional[str] = None) -> None:
        """플레이스홀더에 단일 단락 텍스트 채우기. 여백·줄바꿈으로 겹침 방지."""
        if ph is None:
            return
        try:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            margin_emu = int(0.02 * 914400)
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin_emu
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size_pt)
            p.font.bold = bold
            if font_name:
                p.font.name = font_name
            if color:
                p.font.color.rgb = color
        except Exception as e:
            logger.warning("플레이스홀더 텍스트 채우기 실패: {}", e)

    def _get_layout(self, name: str, fallback_idx: int = 0):
        """
        이름으로 레이아웃 검색 (모든 슬라이드 마스터 탐색).
        guide_template.pptx는 Master[21]에 레이아웃이 있어 이 메서드 필요.
        없으면 prs.slide_layouts[fallback_idx] 사용.

        주요 레이아웃 이름 (guide_template.pptx Master[21]):
          "1_cover_02"  → 표지 슬라이드
          "3_cover_02"  → 섹션 구분 (다크 왼쪽 패널)
          "Title and Content_02"            → 일반 콘텐츠
          "Title and Content_02 (타이틀2줄)" → 2줄 제목 콘텐츠
          "agenda (Print)"                  → 목차
          "last page"                       → 마지막 페이지
        """
        layout = self.template_manager.get_slide_layout_by_name(self.prs, name)
        if layout is not None:
            return layout
        n = len(self.prs.slide_layouts)
        safe_idx = min(fallback_idx, n - 1) if n > 0 else 0
        return self.prs.slide_layouts[safe_idx]

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        is_part_divider: bool = False,
        slogan: Optional[str] = None,
    ) -> None:
        """
        표지 슬라이드 (guide_template 1_cover_02 스타일).
        - 상단 네이비 이미지 영역: (0.36, 0.25) 9.29 × 3.03
        - 제목: (0.26, 3.63) 9.09 × 0.50, 30pt Bold, 네이비
        - 부제목: (0.26, 4.16) 9.09 × 0.30, 15pt
        - 슬로건/태그라인: (0.26, 4.60)
        """
        if is_part_divider:
            # Part 구분자 → section_divider 스타일 재사용
            self.add_section_divider(
                phase_number=0,
                phase_title=_truncate(title, 100),
                phase_subtitle=subtitle or "",
            )
            return

        # 표지 레이아웃 (1_cover_02 우선, 없으면 cover_01, 없으면 첫 레이아웃)
        slide_layout = self._get_layout("1_cover_02") or self._get_layout("cover_01") or self.prs.slide_layouts[0]
        slide = self._add_slide_with_placeholders(slide_layout)

        # 상단 헤더 이미지 영역 (네이비 배경 직사각형)
        # 실측: 1_cover_02 PH17 L=0.356" T=0.253" W=9.290" H=3.031"
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.356), Inches(0.253),
            Inches(9.290), Inches(3.031),
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = STYLE_COLORS["primary"]
        header_rect.line.fill.background()

        # 헤더 내 슬로건 또는 브랜드 태그라인 (선택)
        if slogan and slogan.strip():
            slogan_box = slide.shapes.add_textbox(
                Inches(0.65), Inches(1.30),
                Inches(8.70), Inches(0.80),
            )
            slogan_tf = slogan_box.text_frame
            slogan_tf.word_wrap = True
            slogan_p = slogan_tf.paragraphs[0]
            slogan_p.text = _truncate(slogan.strip(), 120)
            slogan_p.font.name = self.template_manager.get_font_name("title")
            slogan_p.font.size = Pt(20)
            slogan_p.font.bold = False
            slogan_p.font.color.rgb = STYLE_COLORS["sky_blue"]
            slogan_p.alignment = PP_ALIGN.LEFT

        # ── 제목: PH18 사용 (실측 T=3.670" H=0.515"), fallback: textbox T=3.633" ──
        ph18 = self._get_ph(slide, 18)
        if ph18:
            self._fill_ph_text(
                ph18, _truncate(title, 100),
                font_size_pt=30, bold=True,
                color=STYLE_COLORS["primary"],
                font_name=self.template_manager.get_font_name("title"),
            )
        else:
            title_box = slide.shapes.add_textbox(
                Inches(0.258), Inches(3.633),
                Inches(9.087), Inches(0.505),
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            margin_emu = int(0.02 * 914400)
            title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = margin_emu
            title_p = title_tf.paragraphs[0]
            title_p.text = _truncate(title, 100)
            title_p.font.name = self.template_manager.get_font_name("title")
            title_p.font.size = Pt(30)
            title_p.font.bold = True
            title_p.font.color.rgb = STYLE_COLORS["primary"]

        # ── 부제목: PH14 사용 (실측 T=4.201" H=0.227"). 없으면 비워 기본 문구 미노출 ──
        ph14 = self._get_ph(slide, 14)
        if subtitle and str(subtitle).strip():
            if ph14:
                self._fill_ph_text(
                    ph14, _truncate(subtitle, 150),
                    font_size_pt=15, bold=False,
                    color=STYLE_COLORS["secondary"],
                    font_name=self.template_manager.get_font_name("body"),
                )
            else:
                sub_box = slide.shapes.add_textbox(
                    Inches(0.258), Inches(4.164),
                    Inches(9.087), Inches(0.298),
                )
                sub_tf = sub_box.text_frame
                sub_tf.word_wrap = True
                pad_emu = int(0.02 * 914400)
                sub_tf.margin_left = sub_tf.margin_right = sub_tf.margin_top = sub_tf.margin_bottom = pad_emu
                sub_p = sub_tf.paragraphs[0]
                sub_p.text = _truncate(subtitle, 150)
                sub_p.font.name = self.template_manager.get_font_name("body")
                sub_p.font.size = Pt(15)
                sub_p.font.color.rgb = STYLE_COLORS["secondary"]
        elif ph14:
            self._clear_placeholder_text(ph14)

        # 저작권/회사명 (하단 좌측, 실측: Slide 6 Rectangle 1 T=5.080")
        copy_box = slide.shapes.add_textbox(
            Inches(0.258), Inches(5.080),
            Inches(4.339), Inches(0.250),
        )
        copy_tf = copy_box.text_frame
        copy_p = copy_tf.paragraphs[0]
        copy_p.font.name = self.template_manager.get_font_name("body")
        copy_p.font.size = Pt(7)
        copy_p.font.color.rgb = RGBColor(0x6E, 0x83, 0xA1)  # 실측: #6E83A1

    def add_content_slide(
        self,
        title: str,
        bullets: Optional[List[BulletPoint]] = None,
        key_message: Optional[str] = None,
        notes: Optional[str] = None,
        subtitle: Optional[str] = None,
        layout_hint: Optional[str] = None,
    ) -> None:
        """
        콘텐츠 슬라이드 - guide_template.pptx 플레이스홀더 방식 (v4.0).

        Title and Content_02 레이아웃 플레이스홀더 구조 (실측):
          PH13 (BODY)  : T=0.214" - 상단 섹션 라벨 (소형, 비움)
          PH0  (TITLE) : T=0.343" - 슬라이드 제목 (18pt Bold)
          ─────────────────────────── Straight Connector T=0.702" (레이아웃 자동)
          PH14 (BODY)  : T=0.815" - 부제목/헤드라인 (12pt)
          PH15 (BODY)  : T=1.113" - 보조 부제목 (비움)
          PH16 (BODY)  : T=1.333" H=3.779" - 본문 불릿 영역
          ─────────────────────────── Straight Connector T=5.323" (레이아웃 자동)
          PH12 (SLIDE_NUMBER): T=5.362" (레이아웃 자동)
          로고 이미지: 우상단 (레이아웃 자동)
        """
        # 제목 길이에 따라 1줄/2줄 레이아웃 선택
        use_2line = len(title) > 45
        layout_name = "Title and Content_02 (타이틀2줄)" if use_2line else "Title and Content_02"
        slide_layout = self._get_layout(layout_name, fallback_idx=0)
        # 플레이스홀더를 유지한 채 슬라이드 추가 (기본 안내 텍스트만 지움)
        slide = self._add_slide_with_placeholders(slide_layout)

        font_title = self.template_manager.get_font_name("title")
        font_body  = self.template_manager.get_font_name("body")

        # ── PH0: 슬라이드 제목 (18pt Bold, 네이비) ──────────────────
        ph_title = self._get_ph(slide, 0)
        if ph_title:
            self._fill_ph_text(
                ph_title, _truncate(title, MAX_TITLE_CHARS),
                font_size_pt=18, bold=True,
                color=STYLE_COLORS["primary"], font_name=font_title,
            )
        else:
            self._add_content_title(slide, _truncate(title, MAX_TITLE_CHARS))

        # ── PH14: 부제목 (12pt, 블루그레이). 없으면 플레이스홀더 비움(기본 문구 미노출) ──
        ph14 = self._get_ph(slide, 14)
        if subtitle and str(subtitle).strip():
            if ph14:
                self._fill_ph_text(
                    ph14, _truncate(str(subtitle).strip(), 120),
                    font_size_pt=12, bold=False,
                    color=STYLE_COLORS["secondary"], font_name=font_body,
                )
            else:
                self._add_content_subtitle(slide, _truncate(str(subtitle).strip(), 120))
        elif ph14:
            self._clear_placeholder_text(ph14)

        # ── PH16: 본문 불릿 (Body 영역) ──────────────────────────────
        bullets_to_show = (bullets or [])[:MAX_BULLETS_PER_SLIDE]
        if bullets_to_show:
            ph_body = self._get_ph(slide, 16)
            # 타이틀2줄 레이아웃은 PH idx가 4294967295(unlabeled)일 수 있음
            if ph_body is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx >= 4294967000:
                            ph_body = ph
                            break
                    except Exception:
                        pass

            if ph_body:
                tf = ph_body.text_frame
                tf.clear()
                tf.word_wrap = True
                margin_emu = int(0.02 * 914400)
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin_emu
                for i, bullet in enumerate(bullets_to_show):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    bullet_char = "•" if bullet.level == 0 else "–"
                    indent = "    " * min(bullet.level, 2)
                    p.text = f"{indent}{bullet_char} {_truncate(bullet.text, MAX_BULLET_CHARS)}"
                    p.font.name = font_body
                    p.font.size = Pt(12) if bullet.level == 0 else Pt(10)
                    p.font.bold = bullet.emphasis
                    p.font.color.rgb = (
                        STYLE_COLORS["primary"] if bullet.emphasis
                        else self.template_manager.get_color("text_dark")
                    )
                    p.space_after = Pt(4)
            else:
                # 플레이스홀더 없을 때 텍스트박스 fallback (제목/부제목과 겹치지 않음)
                body_box = slide.shapes.add_textbox(
                    Inches(MARGIN_H), Inches(BODY_BOX_TOP),
                    Inches(CONTENT_WIDTH), Inches(BODY_BOX_HEIGHT),
                )
                tf = body_box.text_frame
                tf.word_wrap = True
                margin_emu = int(0.03 * 914400)
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin_emu
                for i, bullet in enumerate(bullets_to_show):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    bullet_char = "•" if bullet.level == 0 else "–"
                    indent = "    " * min(bullet.level, 2)
                    p.text = f"{indent}{bullet_char} {_truncate(bullet.text, MAX_BULLET_CHARS)}"
                    p.font.name = font_body
                    p.font.size = Pt(12) if bullet.level == 0 else Pt(10)
                    p.font.bold = bullet.emphasis
                    p.font.color.rgb = (
                        STYLE_COLORS["primary"] if bullet.emphasis
                        else self.template_manager.get_color("text_dark")
                    )
                    p.space_after = Pt(4)

        # ── 핵심 메시지 바 (하단 강조, 템플릿에 없는 커스텀 요소) ────
        if key_message:
            self._add_key_message(slide, _truncate(key_message, MAX_KEY_MSG_CHARS))

        if notes:
            slide.notes_slide.notes_text_frame.text = (notes or "")[:2000]

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        highlight_rows: Optional[List[int]] = None,
        notes: Optional[str] = None,
    ) -> None:
        """
        테이블 슬라이드 추가

        Args:
            title: 슬라이드 제목
            headers: 테이블 헤더
            rows: 테이블 데이터 행
            highlight_rows: 강조할 행 인덱스
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 제목 추가
        self._add_title_textbox(slide, title)

        # 테이블 생성 (슬라이드 범위 내로 제한)
        rows_count = len(rows) + 1  # 헤더 포함
        cols_count = len(headers)
        if rows_count < 2 or cols_count < 1:
            return
        table_max_height = KEY_MSG_TOP - BODY_BOX_TOP - 0.22
        row_height_inch = 0.26
        table_height = min(row_height_inch * rows_count, table_max_height)
        left = Inches(MARGIN_H)
        top = Inches(BODY_BOX_TOP)
        width = Inches(CONTENT_WIDTH)
        height = Inches(table_height)

        table = slide.shapes.add_table(
            rows_count, cols_count, left, top, width, height
        ).table

        # 열 너비 설정
        col_width = width // cols_count
        for i in range(cols_count):
            table.columns[i].width = col_width

        # 헤더 설정 (길이 제한)
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = _truncate(str(header), MAX_TABLE_CELL_CHARS)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.template_manager.get_color("primary")
            self._format_table_cell(cell, is_header=True)

        # 데이터 행 설정
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= cols_count:
                    break

                cell = table.cell(row_idx + 1, col_idx)
                cell.text = _truncate(str(cell_text) if cell_text else "", MAX_TABLE_CELL_CHARS)

                # 강조 행 처리
                if highlight_rows and row_idx in highlight_rows:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 220)

                self._format_table_cell(cell, is_header=False)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_bullets: List[BulletPoint],
        right_title: str,
        right_bullets: List[BulletPoint],
        notes: Optional[str] = None,
    ) -> None:
        """
        2단 슬라이드 추가

        Args:
            title: 슬라이드 제목
            left_title: 왼쪽 열 제목
            left_bullets: 왼쪽 열 불릿
            right_title: 오른쪽 열 제목
            right_bullets: 오른쪽 열 불릿
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        self._add_title_textbox(slide, _truncate(title, MAX_TITLE_CHARS))

        col_width = (CONTENT_WIDTH - 0.13) / 2
        col_left1 = MARGIN_H
        col_left2 = MARGIN_H + col_width + 0.13
        col_top = BODY_BOX_TOP
        col_height = KEY_MSG_TOP - col_top - 0.05

        left_box = slide.shapes.add_textbox(
            Inches(col_left1), Inches(col_top), Inches(col_width), Inches(col_height)
        )
        self._fill_column(left_box, left_title, left_bullets)

        right_box = slide.shapes.add_textbox(
            Inches(col_left2), Inches(col_top), Inches(col_width), Inches(col_height)
        )
        self._fill_column(right_box, right_title, right_bullets)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def save(self, output_path: Path) -> None:
        """프레젠테이션 저장"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"PPTX 저장 완료: {output_path}")

    def _add_title_textbox(
        self, slide, title: str, size_name: str = "slide_title"
    ) -> None:
        """슬라이드 제목 추가. PH0 플레이스홀더 있으면 사용, 없으면 textbox fallback."""
        ph0 = self._get_ph(slide, 0)
        if ph0:
            self._fill_ph_text(
                ph0, _truncate(title, MAX_TITLE_CHARS),
                font_size_pt=18, bold=True,
                color=STYLE_COLORS["primary"],
                font_name=self.template_manager.get_font_name("title"),
            )
            return
        # fallback: textbox (템플릿 없을 때 겹침 방지용 고정 좌표)
        textbox = slide.shapes.add_textbox(
            Inches(MARGIN_H),
            Inches(TITLE_BOX_TOP),
            Inches(CONTENT_WIDTH),
            Inches(TITLE_BOX_HEIGHT),
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        margin_emu = int(0.02 * 914400)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin_emu
        lines = (title or "").strip().split("\n")
        for i, line in enumerate(lines[:2]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line.strip()
            self._apply_title_format(p, size_name=size_name)

    def _add_content_title(self, slide, title: str) -> None:
        """콘텐츠 슬라이드 제목 바 (18pt Bold, 네이비). 고정 좌표로 겹침 방지."""
        textbox = slide.shapes.add_textbox(
            Inches(MARGIN_H), Inches(TITLE_BOX_TOP),
            Inches(CONTENT_WIDTH), Inches(TITLE_BOX_HEIGHT),
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        pad_emu = int(0.03 * 914400)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = pad_emu
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.template_manager.get_font_name("title")
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = STYLE_COLORS["primary"]

    def _add_content_subtitle(self, slide, subtitle: str) -> None:
        """콘텐츠 슬라이드 부제목 바 (15pt). 제목 영역과 겹치지 않도록 고정 좌표."""
        textbox = slide.shapes.add_textbox(
            Inches(MARGIN_H), Inches(SUBTITLE_BOX_TOP),
            Inches(CONTENT_WIDTH), Inches(SUBTITLE_BOX_HEIGHT),
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        pad_emu = int(0.03 * 914400)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = pad_emu
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = self.template_manager.get_font_name("body")
        p.font.size = Pt(15)
        p.font.bold = False
        p.font.color.rgb = STYLE_COLORS["secondary"]

    def _add_deco_line(self, slide) -> None:
        """장식 수평 라인 (guide_template: y=1.34, 네이비, 얇은 사각형)."""
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_H), Inches(DECO_LINE_TOP),
            Inches(CONTENT_WIDTH), Inches(0.02),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = STYLE_COLORS["primary"]
        line_shape.line.fill.background()

    def _add_key_message(self, slide, message: str) -> None:
        """슬라이드 하단 핵심 메시지 바 (네이비 배경, 박스 안 글자 잘림 방지)."""
        if not (message or "").strip():
            return
        # 배경 바 (네이비)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_H), Inches(KEY_MSG_TOP),
            Inches(CONTENT_WIDTH), Inches(KEY_MSG_HEIGHT),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS["primary"]
        bg.line.fill.background()
        # 텍스트 (내부 여백 확대, 높이 여유로 2~3줄 수납)
        pad_h, pad_v = 0.22, 0.10
        textbox = slide.shapes.add_textbox(
            Inches(MARGIN_H + pad_h), Inches(KEY_MSG_TOP + pad_v),
            Inches(CONTENT_WIDTH - 2 * pad_h), Inches(KEY_MSG_HEIGHT - 2 * pad_v),
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        margin_emu = int(0.03 * 914400)
        tf.margin_left = margin_emu
        tf.margin_right = margin_emu
        tf.margin_top = margin_emu
        tf.margin_bottom = margin_emu
        p = tf.paragraphs[0]
        p.text = (message or "").strip()
        p.font.name = self.template_manager.get_font_name("body")
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = STYLE_COLORS["white"]
        p.alignment = PP_ALIGN.LEFT

    def _apply_title_format(self, paragraph, size_name: str = "slide_title") -> None:
        """제목 포맷 적용"""
        paragraph.font.size = self.template_manager.get_font_size(size_name)
        paragraph.font.bold = True
        paragraph.font.name = self.template_manager.get_font_name("title")
        paragraph.font.color.rgb = self.template_manager.get_color("primary")

    def _apply_text_format(
        self, paragraph, size_name: str = "body", color_name: str = "text_dark"
    ) -> None:
        """텍스트 포맷 적용"""
        paragraph.font.size = self.template_manager.get_font_size(size_name)
        paragraph.font.name = self.template_manager.get_font_name("body")
        paragraph.font.color.rgb = self.template_manager.get_color(color_name)

    def _format_table_cell(self, cell, is_header: bool = False) -> None:
        """테이블 셀 포맷 적용 (줄바꿈·작은 폰트·내부 여백으로 셀 안에 수납)"""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        margin_emu = int(0.02 * 914400)
        tf.margin_left = margin_emu
        tf.margin_right = margin_emu
        tf.margin_top = margin_emu
        tf.margin_bottom = margin_emu

        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(9) if not is_header else Pt(10)
            paragraph.font.name = self.template_manager.get_font_name("body")
            paragraph.font.bold = is_header
            paragraph.alignment = PP_ALIGN.CENTER

            if is_header:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            else:
                paragraph.font.color.rgb = self.template_manager.get_color("text_dark")

    def _fill_column(
        self, textbox, column_title: str, bullets: List[BulletPoint]
    ) -> None:
        """2단 레이아웃의 열 채우기 (word_wrap·여백으로 겹침 방지)"""
        tf = textbox.text_frame
        tf.word_wrap = True
        margin_emu = int(0.02 * 914400)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin_emu

        p = tf.paragraphs[0]
        p.text = _truncate(column_title, MAX_TITLE_CHARS)
        p.font.name = self.template_manager.get_font_name("title")
        p.font.size = self.template_manager.get_font_size("subtitle")
        p.font.bold = True
        p.font.color.rgb = self.template_manager.get_color("secondary")

        for bullet in bullets[:MAX_BULLETS_PER_SLIDE]:
            p = tf.add_paragraph()
            p.text = _truncate(bullet.text, MAX_BULLET_CHARS)
            p.level = min(bullet.level, 2)
            p.font.size = self.template_manager.get_font_size("body")
            p.font.name = self.template_manager.get_font_name("body")
            p.font.bold = bullet.emphasis

    # ==========================================================================
    # 추가 레이아웃
    # ==========================================================================
    def add_three_column_slide(
        self,
        title: str,
        columns: List[dict],
        notes: Optional[str] = None,
    ) -> None:
        """
        3단 레이아웃 슬라이드 추가

        Args:
            title: 슬라이드 제목
            columns: [{"title": "열1", "content": "내용", "icon": "★"}, ...]
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        self._add_title_textbox(slide, _truncate(title, MAX_TITLE_CHARS))

        col_width = (CONTENT_WIDTH - 0.40) / 3  # 3열 균등 배분
        col_height = min(3.0, KEY_MSG_TOP - BODY_BOX_TOP - 0.1)
        left_start = MARGIN_H
        top = BODY_BOX_TOP
        gap = 0.20

        for i, col in enumerate(columns[:3]):
            left = left_start + i * (col_width + gap)
            self._add_column_box(slide, col, left, top, col_width, col_height)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_column_box(
        self,
        slide,
        column: dict,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """3단 레이아웃의 열 박스"""
        # 배경 박스
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        shape.line.color.rgb = self.template_manager.get_color("primary")
        shape.line.width = Pt(1.5)

        content_left = left + 0.2
        content_width = width - 0.4
        current_top = top + 0.2

        # 아이콘 (있는 경우)
        if column.get("icon"):
            icon_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top),
                Inches(content_width),
                Inches(0.6),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = column["icon"]
            icon_p.font.name = self.template_manager.get_font_name("body")
            icon_p.font.size = Pt(28)
            icon_p.alignment = PP_ALIGN.CENTER
            current_top += 0.6

        # 열 제목
        title_box = slide.shapes.add_textbox(
            Inches(content_left),
            Inches(current_top),
            Inches(content_width),
            Inches(0.5),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = _truncate(column.get("title", ""), 50)
        title_p.font.name = self.template_manager.get_font_name("title")
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = self.template_manager.get_color("primary")
        title_p.alignment = PP_ALIGN.CENTER

        if column.get("content"):
            content_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top + 0.6),
                Inches(content_width),
                Inches(height - current_top - 0.5),
            )
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_p = content_tf.paragraphs[0]
            content_p.text = _truncate(column["content"], 250)
            content_p.font.name = self.template_manager.get_font_name("body")
            content_p.font.size = Pt(11)
            content_p.font.color.rgb = self.template_manager.get_color("text_dark")

        # 불릿 (있는 경우)
        if column.get("bullets"):
            bullets_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(current_top + 0.6),
                Inches(content_width),
                Inches(height - current_top - 0.5),
            )
            bullets_tf = bullets_box.text_frame
            bullets_tf.word_wrap = True

            for j, bullet in enumerate(column["bullets"][:8]):
                if j == 0:
                    p = bullets_tf.paragraphs[0]
                else:
                    p = bullets_tf.add_paragraph()
                p.text = "• " + _truncate(str(bullet), 120)
                p.font.name = self.template_manager.get_font_name("body")
                p.font.size = Pt(10)

    def add_big_number_slide(
        self,
        title: str,
        stats: List[dict],
        notes: Optional[str] = None,
    ) -> None:
        """
        큰 숫자 강조 슬라이드 (KPI, 성과 등)

        Args:
            title: 슬라이드 제목
            stats: [{"value": "95%", "label": "만족도", "description": "목표 대비 +5%"}, ...]
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 통계 카드들
        num_stats = min(len(stats), 4)
        total_width = CONTENT_WIDTH - 0.10
        card_width = total_width / num_stats
        left_start = MARGIN_H
        top = BODY_BOX_TOP

        for i, stat in enumerate(stats[:4]):
            left = left_start + i * card_width
            self._add_stat_card(slide, stat, left, top, card_width - 0.3)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_stat_card(
        self,
        slide,
        stat: dict,
        left: float,
        top: float,
        width: float,
    ) -> None:
        """통계 카드 (word_wrap, 길이 제한)"""
        value_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.2),
        )
        value_tf = value_box.text_frame
        value_tf.word_wrap = True
        value_p = value_tf.paragraphs[0]
        value_p.text = _truncate(stat.get("value", ""), 30)
        value_p.font.name = self.template_manager.get_font_name("title")
        value_p.font.size = Pt(40)
        value_p.font.bold = True
        value_p.font.color.rgb = self.template_manager.get_color("primary")
        value_p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top + 1.2),
            Inches(width),
            Inches(0.5),
        )
        label_tf = label_box.text_frame
        label_tf.word_wrap = True
        label_p = label_tf.paragraphs[0]
        label_p.text = _truncate(stat.get("label", ""), 40)
        label_p.font.name = self.template_manager.get_font_name("body")
        label_p.font.size = Pt(14)
        label_p.font.bold = True
        label_p.font.color.rgb = self.template_manager.get_color("text_dark")
        label_p.alignment = PP_ALIGN.CENTER

        if stat.get("description"):
            desc_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top + 1.7),
                Inches(width),
                Inches(0.6),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = _truncate(stat["description"], 80)
            desc_p.font.name = self.template_manager.get_font_name("body")
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = self.template_manager.get_color("text_light")
            desc_p.alignment = PP_ALIGN.CENTER

    def add_icon_grid_slide(
        self,
        title: str,
        items: List[dict],
        columns: int = 4,
        notes: Optional[str] = None,
    ) -> None:
        """
        아이콘 그리드 슬라이드

        Args:
            title: 슬라이드 제목
            items: [{"icon": "🎯", "title": "목표", "description": "설명"}, ...]
            columns: 열 개수 (3 또는 4)
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 그리드 아이템
        total_width = CONTENT_WIDTH - 0.10
        margin = 0.18
        item_width = (total_width - margin * (columns + 1)) / columns
        item_height = 1.35
        left_start = MARGIN_H
        top = BODY_BOX_TOP

        for i, item in enumerate(items[:columns * 2]):  # 최대 2행
            row = i // columns
            col = i % columns

            item_left = left_start + margin + col * (item_width + margin)
            item_top = top + row * (item_height + margin)

            self._add_icon_item(slide, item, item_left, item_top, item_width, item_height)

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_icon_item(
        self,
        slide,
        item: dict,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """아이콘 그리드 아이템"""
        # 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
        shape.line.color.rgb = self.template_manager.get_color("text_light")
        shape.line.width = Pt(1)

        # 아이콘
        if item.get("icon"):
            icon_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 0.1),
                Inches(width - 0.2),
                Inches(0.5),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = item["icon"]
            icon_p.font.name = self.template_manager.get_font_name("body")
            icon_p.font.size = Pt(24)
            icon_p.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.1),
            Inches(top + 0.6),
            Inches(width - 0.2),
            Inches(0.4),
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = item.get("title", "")
        title_p.font.name = self.template_manager.get_font_name("title")
        title_p.font.size = Pt(11)
        title_p.font.bold = True
        title_p.font.color.rgb = self.template_manager.get_color("primary")
        title_p.alignment = PP_ALIGN.CENTER

        # 설명
        if item.get("description"):
            desc_box = slide.shapes.add_textbox(
                Inches(left + 0.1),
                Inches(top + 1.0),
                Inches(width - 0.2),
                Inches(height - 1.1),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = item["description"]
            desc_p.font.name = self.template_manager.get_font_name("body")
            desc_p.font.size = Pt(9)
            desc_p.font.color.rgb = self.template_manager.get_color("text_dark")
            desc_p.alignment = PP_ALIGN.CENTER

    def add_quote_slide(
        self,
        title: str,
        quote: str,
        author: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> None:
        """
        인용문/추천사 슬라이드

        Args:
            title: 슬라이드 제목
            quote: 인용문 내용
            author: 작성자 정보
            notes: 발표자 노트
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 큰 따옴표 배경 (슬라이드 내 범위)
        quote_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN_H + 0.30),
            Inches(BODY_BOX_TOP),
            Inches(CONTENT_WIDTH - 0.60),
            Inches(2.40),
        )
        quote_bg.fill.solid()
        quote_bg.fill.fore_color.rgb = self.template_manager.get_color("background_light")
        quote_bg.line.fill.background()

        # 인용문 (길이 제한, word_wrap)
        quote_display = _truncate(quote, 400)
        quote_box = slide.shapes.add_textbox(
            Inches(MARGIN_H + 0.40),
            Inches(BODY_BOX_TOP + 0.15),
            Inches(CONTENT_WIDTH - 0.80),
            Inches(1.80),
        )
        quote_tf = quote_box.text_frame
        quote_tf.word_wrap = True
        quote_p = quote_tf.paragraphs[0]
        quote_p.text = f'"{quote_display}"'
        quote_p.font.name = self.template_manager.get_font_name("body")
        quote_p.font.size = Pt(18)
        quote_p.font.italic = True
        quote_p.font.color.rgb = self.template_manager.get_color("text_dark")
        quote_p.alignment = PP_ALIGN.CENTER

        # 작성자
        if author:
            author_display = _truncate(author, 80)
            author_box = slide.shapes.add_textbox(
                Inches(MARGIN_H),
                Inches(4.6),
                Inches(CONTENT_WIDTH),
                Inches(0.5),
            )
            author_tf = author_box.text_frame
            author_tf.word_wrap = True
            author_p = author_tf.paragraphs[0]
            author_p.text = f"- {author_display}"
            author_p.font.name = self.template_manager.get_font_name("body")
            author_p.font.size = Pt(14)
            author_p.font.color.rgb = self.template_manager.get_color("text_light")
            author_p.alignment = PP_ALIGN.RIGHT

        # 발표자 노트
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # ========================================
    # Modern 스타일 슬라이드 메서드
    # ========================================

    def add_teaser_slide(
        self,
        headline: str,
        subheadline: str = "",
        background_color: str = "dark_bg",
        notes: str = "",
    ):
        """
        티저/HOOK 슬라이드 (guide_template 기준 크기, 다크 배경).
        """
        slide_layout = self._get_layout("3_cover_02", fallback_idx=0)
        slide = self._add_slide_with_placeholders(slide_layout)

        # 다크 배경 (슬라이드 전체)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(SLIDE_WIDTH_INCHES), Inches(SLIDE_HEIGHT_INCHES),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS.get(background_color, STYLE_COLORS["dark_bg"])
        bg.line.fill.background()

        # 메인 헤드라인
        headline_display = _truncate(headline, 120)
        headline_box = slide.shapes.add_textbox(
            Inches(MARGIN_H), Inches(1.60),
            Inches(CONTENT_WIDTH), Inches(1.50),
        )
        headline_tf = headline_box.text_frame
        headline_tf.word_wrap = True
        headline_p = headline_tf.paragraphs[0]
        headline_p.text = headline_display
        headline_p.font.name = self.template_manager.get_font_name("title")
        headline_p.font.size = Pt(30)
        headline_p.font.bold = True
        headline_p.font.color.rgb = STYLE_COLORS["white"]
        headline_p.alignment = PP_ALIGN.CENTER

        # 서브 헤드라인
        if subheadline:
            sub_display = _truncate(subheadline, 200)
            sub_box = slide.shapes.add_textbox(
                Inches(MARGIN_H), Inches(3.20),
                Inches(CONTENT_WIDTH), Inches(0.90),
            )
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = sub_display
            sub_p.font.name = self.template_manager.get_font_name("body")
            sub_p.font.size = Pt(15)
            sub_p.font.color.rgb = STYLE_COLORS["sky_blue"]
            sub_p.alignment = PP_ALIGN.CENTER

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_section_divider(
        self,
        phase_number: int,
        phase_title: str,
        phase_subtitle: str = "",
        notes: str = "",
    ):
        """
        섹션 구분 슬라이드 (guide_template 3_cover_02 스타일).
        - 왼쪽 다크 네이비 패널: (0, 0) 4.43 × 5.625
        - 섹션 번호: (0.49, 1.99) 15pt, 라이트 블루
        - 섹션 제목: (0.49, 2.37) 3.64 × 1.30, 30pt Bold, 흰색
        - 오른쪽: 흰색 배경 + 장식 요소
        """
        slide_layout = self._get_layout("3_cover_02", fallback_idx=0)
        slide = self._add_slide_with_placeholders(slide_layout)

        # ── 실측 좌표 (Slide 14 기준) ─────────────────────────────────
        # 직사각형: L=0.000" T=0.000" W=4.427" H=5.625" (다크 네이비 패널)
        # 번호:     L=0.495" T=1.992" W=1.017" H=0.302" | 15pt
        # 제목:     L=0.495" T=2.366" W=3.645" H=1.305" | 30pt Bold
        LEFT_PANEL_W  = 4.427
        PANEL_PAD_LEFT = 0.495

        # 왼쪽 다크 네이비 패널 (전체 높이)
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(LEFT_PANEL_W), Inches(SLIDE_HEIGHT_INCHES),
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = STYLE_COLORS["section_dark"]
        left_panel.line.fill.background()

        # 섹션 번호 (예: "01", "02" ...) - 실측: T=1.992" 15pt
        num_str = f"{phase_number:02d}" if phase_number < 100 else str(phase_number)
        num_box = slide.shapes.add_textbox(
            Inches(PANEL_PAD_LEFT), Inches(1.992),
            Inches(1.017), Inches(0.302),
        )
        num_tf = num_box.text_frame
        pad_emu = int(0.02 * 914400)
        num_tf.margin_left = num_tf.margin_right = num_tf.margin_top = num_tf.margin_bottom = pad_emu
        num_p = num_tf.paragraphs[0]
        num_p.text = num_str
        num_p.font.name = self.template_manager.get_font_name("title")
        num_p.font.size = Pt(15)
        num_p.font.bold = False
        num_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 섹션 제목 - 실측: T=2.366" W=3.645" H=1.305" | 30pt Bold 흰색
        title_display = _truncate(phase_title, 55)
        title_box = slide.shapes.add_textbox(
            Inches(PANEL_PAD_LEFT), Inches(2.366),
            Inches(3.645), Inches(1.305),
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = pad_emu
        title_p = title_tf.paragraphs[0]
        title_p.text = title_display
        title_p.font.name = self.template_manager.get_font_name("title")
        title_p.font.size = Pt(30)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["white"]
        title_p.space_after = Pt(0)

        # 섹션 부제목 (있으면)
        if phase_subtitle:
            sub_display = _truncate(phase_subtitle, 90)
            sub_box = slide.shapes.add_textbox(
                Inches(PANEL_PAD_LEFT), Inches(3.80),
                Inches(3.645), Inches(0.70),
            )
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_tf.margin_left = sub_tf.margin_right = sub_tf.margin_top = sub_tf.margin_bottom = pad_emu
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = sub_display
            sub_p.font.name = self.template_manager.get_font_name("body")
            sub_p.font.size = Pt(12)
            sub_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 오른쪽 패널: 흰색 배경
        right_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(LEFT_PANEL_W), Inches(0),
            Inches(SLIDE_WIDTH_INCHES - LEFT_PANEL_W), Inches(SLIDE_HEIGHT_INCHES),
        )
        right_panel.fill.solid()
        right_panel.fill.fore_color.rgb = STYLE_COLORS["white"]
        right_panel.line.fill.background()

        # 오른쪽 장식 라인 (스카이 블루)
        right_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(LEFT_PANEL_W + 0.40), Inches(2.80),
            Inches(SLIDE_WIDTH_INCHES - LEFT_PANEL_W - 0.80), Inches(0.02),
        )
        right_line.fill.solid()
        right_line.fill.fore_color.rgb = STYLE_COLORS["sky_blue"]
        right_line.line.fill.background()

        # 발표자 노트
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_key_message_slide(
        self,
        message: str,
        supporting_text: str = "",
        background_style: str = "gradient",
        notes: str = "",
    ):
        """
        핵심 메시지 슬라이드 - 중앙 정렬된 임팩트 메시지
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 배경 (다크 또는 라이트)
        if background_style == "dark":
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(SLIDE_WIDTH_INCHES), Inches(SLIDE_HEIGHT_INCHES),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = STYLE_COLORS["dark_bg"]
            bg.line.fill.background()
            text_color = STYLE_COLORS["white"]
        else:
            text_color = STYLE_COLORS["primary"]

        # 핵심 메시지 (guide_template 기준 크기)
        msg_display = _truncate(message, 200)
        msg_box = slide.shapes.add_textbox(
            Inches(MARGIN_H), Inches(1.60),
            Inches(CONTENT_WIDTH), Inches(1.80),
        )
        msg_tf = msg_box.text_frame
        msg_tf.word_wrap = True
        msg_p = msg_tf.paragraphs[0]
        msg_p.text = msg_display
        msg_p.font.name = self.template_manager.get_font_name("title")
        msg_p.font.size = Pt(24)
        msg_p.font.bold = True
        msg_p.font.color.rgb = text_color
        msg_p.alignment = PP_ALIGN.CENTER

        # 보조 텍스트
        if supporting_text:
            sup_display = _truncate(supporting_text, 300)
            sup_box = slide.shapes.add_textbox(
                Inches(MARGIN_H + 0.30), Inches(3.50),
                Inches(CONTENT_WIDTH - 0.60), Inches(1.00),
            )
            sup_tf = sup_box.text_frame
            sup_tf.word_wrap = True
            sup_p = sup_tf.paragraphs[0]
            sup_p.text = sup_display
            sup_p.font.name = self.template_manager.get_font_name("body")
            sup_p.font.size = Pt(14)
            sup_p.font.color.rgb = STYLE_COLORS["sky_blue"] if background_style == "dark" else STYLE_COLORS["text_gray"]
            sup_p.alignment = PP_ALIGN.CENTER

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_comparison_slide(
        self,
        title: str,
        as_is: dict,
        to_be: dict,
        notes: str = "",
    ):
        """
        AS-IS / TO-BE 비교 슬라이드
        Modern 스타일: 좌우 분할, 시각적 대비
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        self._add_title_textbox(slide, _truncate(title, MAX_TITLE_CHARS))

        col_w = (CONTENT_WIDTH - 0.33) / 2
        col_h = KEY_MSG_TOP - BODY_BOX_TOP - 0.1

        # AS-IS 영역 (왼쪽)
        as_is_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN_H),
            Inches(BODY_BOX_TOP),
            Inches(col_w),
            Inches(col_h),
        )
        as_is_bg.fill.solid()
        as_is_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        as_is_bg.line.fill.background()

        as_title_box = slide.shapes.add_textbox(
            Inches(MARGIN_H + 0.2),
            Inches(BODY_BOX_TOP + 0.15),
            Inches(col_w - 0.4),
            Inches(0.5),
        )
        as_title_tf = as_title_box.text_frame
        as_title_tf.word_wrap = True
        as_title_p = as_title_tf.paragraphs[0]
        as_title_p.text = _truncate(as_is.get("title", "AS-IS (현재)"), 40)
        as_title_p.font.name = self.template_manager.get_font_name("title")
        as_title_p.font.size = Pt(20)
        as_title_p.font.bold = True
        as_title_p.font.color.rgb = STYLE_COLORS["text_gray"]

        as_content_box = slide.shapes.add_textbox(
            Inches(MARGIN_H + 0.2),
            Inches(BODY_BOX_TOP + 0.7),
            Inches(col_w - 0.4),
            Inches(col_h - 0.85),
        )
        as_content_tf = as_content_box.text_frame
        as_content_tf.word_wrap = True
        as_items = list(as_is.get("items", []))[:MAX_BULLETS_PER_SLIDE]
        for i, item in enumerate(as_items):
            p = as_content_tf.paragraphs[0] if i == 0 else as_content_tf.add_paragraph()
            p.text = "• " + _truncate(str(item), 180)
            p.font.name = self.template_manager.get_font_name("body")
            p.font.size = Pt(12)
            p.font.color.rgb = STYLE_COLORS["text_gray"]
            p.space_after = Pt(6)

        # TO-BE 영역 (오른쪽)
        to_left = MARGIN_H + col_w + 0.33
        to_be_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(to_left),
            Inches(BODY_BOX_TOP),
            Inches(col_w),
            Inches(col_h),
        )
        to_be_bg.fill.solid()
        to_be_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        to_be_bg.line.fill.background()

        to_title_box = slide.shapes.add_textbox(
            Inches(to_left + 0.2),
            Inches(BODY_BOX_TOP + 0.15),
            Inches(col_w - 0.4),
            Inches(0.5),
        )
        to_title_tf = to_title_box.text_frame
        to_title_tf.word_wrap = True
        to_title_p = to_title_tf.paragraphs[0]
        to_title_p.text = _truncate(to_be.get("title", "TO-BE (제안)"), 40)
        to_title_p.font.name = self.template_manager.get_font_name("title")
        to_title_p.font.size = Pt(20)
        to_title_p.font.bold = True
        to_title_p.font.color.rgb = RGBColor(255, 255, 255)

        to_content_box = slide.shapes.add_textbox(
            Inches(to_left + 0.2),
            Inches(BODY_BOX_TOP + 0.7),
            Inches(col_w - 0.4),
            Inches(col_h - 0.85),
        )
        to_content_tf = to_content_box.text_frame
        to_content_tf.word_wrap = True
        to_items = list(to_be.get("items", []))[:MAX_BULLETS_PER_SLIDE]
        for i, item in enumerate(to_items):
            p = to_content_tf.paragraphs[0] if i == 0 else to_content_tf.add_paragraph()
            p.text = "• " + _truncate(str(item), 180)
            p.font.name = self.template_manager.get_font_name("body")
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(6)

        # 화살표 (중앙 간격)
        arrow_left = MARGIN_H + col_w + 0.33 / 2 - 0.2
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(arrow_left),
            Inches(BODY_BOX_TOP + col_h / 2 - 0.25),
            Inches(0.4),
            Inches(0.35),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = STYLE_COLORS["sky_blue"]
        arrow.line.fill.background()

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_index_slide(
        self,
        title: str = "목차",
        items: list = None,
        current_index: int = -1,
        notes: str = "",
    ):
        """
        목차/인덱스 슬라이드 (guide_template "agenda (Print)" 스타일).
        - 좌측 제목 헤더 영역
        - 우측 항목 리스트 (번호 + 텍스트)
        """
        items = items or []
        slide_layout = self._get_layout("agenda (Print)") or self._get_layout("agenda") or self._get_layout("Title and Content_02", 0)
        slide = self._add_slide_with_placeholders(slide_layout)

        # 좌측 타이틀 영역 (네이비 배경)
        left_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(3.20), Inches(SLIDE_HEIGHT_INCHES),
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = STYLE_COLORS["primary"]
        left_bg.line.fill.background()

        # 목차 타이틀 텍스트 (흰색, 좌측 여백으로 몰림 방지)
        toc_left, toc_top, toc_w, toc_h = 0.45, 1.75, 2.45, 1.25
        toc_title_box = slide.shapes.add_textbox(
            Inches(toc_left), Inches(toc_top),
            Inches(toc_w), Inches(toc_h),
        )
        toc_title_tf = toc_title_box.text_frame
        toc_title_tf.word_wrap = True
        toc_pad = int(0.05 * 914400)
        toc_title_tf.margin_left = toc_pad
        toc_title_tf.margin_right = toc_pad
        toc_title_tf.margin_top = toc_pad
        toc_title_tf.margin_bottom = toc_pad
        toc_title_p = toc_title_tf.paragraphs[0]
        toc_title_p.text = _truncate(title, 60)
        toc_title_p.font.name = self.template_manager.get_font_name("title")
        toc_title_p.font.size = Pt(24)
        toc_title_p.font.bold = True
        toc_title_p.font.color.rgb = STYLE_COLORS["white"]

        # 목차 항목들 (우측, 항목-구분선 간격 확대)
        start_y = 0.45
        row_h = min(0.58, (SLIDE_HEIGHT_INCHES - start_y - 0.35) / max(len(items), 1))
        for i, item in enumerate(items[:8]):
            is_current = (i == current_index)
            item_top = start_y + i * row_h

            # 번호 원형 배경
            num_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(3.45), Inches(item_top + 0.06),
                Inches(0.38), Inches(0.38),
            )
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = STYLE_COLORS["accent"] if is_current else STYLE_COLORS["secondary"]
            num_bg.line.fill.background()

            # 번호
            num_box = slide.shapes.add_textbox(
                Inches(3.45), Inches(item_top + 0.08),
                Inches(0.38), Inches(0.34),
            )
            num_p = num_box.text_frame.paragraphs[0]
            num_p.text = f"{i + 1:02d}"
            num_p.font.size = Pt(10)
            num_p.font.bold = True
            num_p.font.color.rgb = STYLE_COLORS["white"]
            num_p.alignment = PP_ALIGN.CENTER

            # 항목 텍스트 (상단 글자-하단 구분선 간격 확보)
            item_box = slide.shapes.add_textbox(
                Inches(3.93), Inches(item_top + 0.06),
                Inches(5.52), Inches(row_h - 0.08),
            )
            item_tf = item_box.text_frame
            item_tf.word_wrap = True
            item_p = item_tf.paragraphs[0]
            item_p.text = _truncate(str(item), 70)
            item_p.font.size = Pt(14) if is_current else Pt(12)
            item_p.font.bold = is_current
            item_p.font.name = self.template_manager.get_font_name("body")
            item_p.font.color.rgb = STYLE_COLORS["primary"] if is_current else STYLE_COLORS["text_dark"]

            # 구분선 (각 항목 하단, 글자와 간격 유지)
            if i < len(items) - 1:
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(3.45), Inches(item_top + row_h - 0.04),
                    Inches(6.00), Inches(0.01),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = STYLE_COLORS["light"]
                sep.line.fill.background()

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_content_example_slide(
        self,
        title: str,
        examples: list,
        notes: str = "",
    ):
        """
        콘텐츠 예시 슬라이드 - 마케팅/PR용
        Modern 스타일: 카드형 레이아웃, 이미지 플레이스홀더
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 최대 3개 카드
        num_cards = min(len(examples), 3)
        card_width = 3.8
        gap = 0.3
        total_width = num_cards * card_width + (num_cards - 1) * gap
        start_x = (SLIDE_WIDTH_INCHES - total_width) / 2

        for i, example in enumerate(examples[:3]):
            x = start_x + i * (card_width + gap)

            # 카드 배경
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(1.6),
                Inches(card_width),
                Inches(5.5),
            )
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            card_bg.line.color.rgb = RGBColor(230, 230, 230)

            # 이미지 플레이스홀더
            img_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.15),
                Inches(1.75),
                Inches(card_width - 0.3),
                Inches(2.5),
            )
            img_placeholder.fill.solid()
            img_placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
            img_placeholder.line.fill.background()

            # 이미지 아이콘 텍스트
            icon_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(2.7),
                Inches(card_width - 0.3),
                Inches(0.5),
            )
            icon_tf = icon_box.text_frame
            icon_p = icon_tf.paragraphs[0]
            icon_p.text = "📷 이미지"
            icon_p.font.size = Pt(14)
            icon_p.font.color.rgb = RGBColor(150, 150, 150)
            icon_p.alignment = PP_ALIGN.CENTER

            # 콘텐츠 유형
            type_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(4.35),
                Inches(card_width - 0.3),
                Inches(0.4),
            )
            type_tf = type_box.text_frame
            type_p = type_tf.paragraphs[0]
            content_type = example.get("content_type", "콘텐츠")
            type_p.text = content_type
            type_p.font.size = Pt(11)
            type_p.font.bold = True
            type_p.font.color.rgb = STYLE_COLORS["sky_blue"]

            # 제목
            ex_title_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(4.7),
                Inches(card_width - 0.3),
                Inches(0.6),
            )
            ex_title_tf = ex_title_box.text_frame
            ex_title_tf.word_wrap = True
            ex_title_p = ex_title_tf.paragraphs[0]
            ex_title_p.text = example.get("title", "")
            ex_title_p.font.size = Pt(13)
            ex_title_p.font.bold = True
            ex_title_p.font.color.rgb = STYLE_COLORS["dark_blue"]

            # 설명
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(5.3),
                Inches(card_width - 0.3),
                Inches(1.0),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = example.get("description", "")[:80]
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = STYLE_COLORS["text_gray"]

            # 채널/해시태그
            channel_box = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(KEY_MSG_TOP - 0.50),
                Inches(card_width - 0.3),
                Inches(0.4),
            )
            channel_tf = channel_box.text_frame
            channel_p = channel_tf.paragraphs[0]
            channel = example.get("channel", "")
            channel_p.text = f"#{channel}" if channel else ""
            channel_p.font.size = Pt(9)
            channel_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_channel_strategy_slide(
        self,
        title: str,
        channels: list,
        notes: str = "",
    ):
        """
        채널 전략 슬라이드 - 채널별 역할/KPI
        Modern 스타일: 채널 아이콘, 역할, KPI 표시
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 채널별 영역 (최대 4개)
        num_channels = min(len(channels), 4)
        col_width = 2.9
        gap = 0.3
        total_width = num_channels * col_width + (num_channels - 1) * gap
        start_x = (SLIDE_WIDTH_INCHES - total_width) / 2

        channel_colors = [
            STYLE_COLORS["dark_blue"],
            STYLE_COLORS["sky_blue"],
            RGBColor(230, 126, 34),  # Orange
            RGBColor(155, 89, 182),  # Purple
        ]

        for i, channel in enumerate(channels[:4]):
            x = start_x + i * (col_width + gap)
            color = channel_colors[i % len(channel_colors)]

            # 채널 헤더
            header_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(1.6),
                Inches(col_width),
                Inches(0.8),
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = color
            header_bg.line.fill.background()

            # 채널명
            name_box = slide.shapes.add_textbox(
                Inches(x),
                Inches(1.75),
                Inches(col_width),
                Inches(0.5),
            )
            name_tf = name_box.text_frame
            name_p = name_tf.paragraphs[0]
            name_p.text = channel.get("name", "")
            name_p.font.size = Pt(16)
            name_p.font.bold = True
            name_p.font.color.rgb = RGBColor(255, 255, 255)
            name_p.alignment = PP_ALIGN.CENTER

            # 역할 영역
            role_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(2.4),
                Inches(col_width),
                Inches(2.2),
            )
            role_bg.fill.solid()
            role_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            role_bg.line.color.rgb = RGBColor(230, 230, 230)

            # 역할 제목
            role_title_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(2.5),
                Inches(col_width - 0.2),
                Inches(0.4),
            )
            role_title_tf = role_title_box.text_frame
            role_title_p = role_title_tf.paragraphs[0]
            role_title_p.text = "역할"
            role_title_p.font.size = Pt(11)
            role_title_p.font.bold = True
            role_title_p.font.color.rgb = color

            # 역할 내용
            role_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(2.9),
                Inches(col_width - 0.2),
                Inches(1.5),
            )
            role_tf = role_box.text_frame
            role_tf.word_wrap = True
            role_p = role_tf.paragraphs[0]
            role_p.text = channel.get("role", "")
            role_p.font.size = Pt(10)
            role_p.font.color.rgb = STYLE_COLORS["text_gray"]

            # KPI 영역
            kpi_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(4.6),
                Inches(col_width),
                Inches(2.4),
            )
            kpi_bg.fill.solid()
            kpi_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
            kpi_bg.line.color.rgb = RGBColor(230, 230, 230)

            # KPI 제목
            kpi_title_box = slide.shapes.add_textbox(
                Inches(x + 0.1),
                Inches(4.7),
                Inches(col_width - 0.2),
                Inches(0.4),
            )
            kpi_title_tf = kpi_title_box.text_frame
            kpi_title_p = kpi_title_tf.paragraphs[0]
            kpi_title_p.text = "KPI"
            kpi_title_p.font.size = Pt(11)
            kpi_title_p.font.bold = True
            kpi_title_p.font.color.rgb = color

            # KPI 항목들
            kpis = channel.get("kpis", [])
            for j, kpi in enumerate(kpis[:3]):
                kpi_item_box = slide.shapes.add_textbox(
                    Inches(x + 0.1),
                    Inches(5.1 + j * 0.6),
                    Inches(col_width - 0.2),
                    Inches(0.55),
                )
                kpi_item_tf = kpi_item_box.text_frame
                kpi_item_tf.word_wrap = True

                # KPI 이름
                kpi_name_p = kpi_item_tf.paragraphs[0]
                kpi_name_p.text = kpi.get("name", "")
                kpi_name_p.font.size = Pt(9)
                kpi_name_p.font.color.rgb = STYLE_COLORS["text_gray"]

                # KPI 값
                kpi_value_p = kpi_item_tf.add_paragraph()
                kpi_value_p.text = kpi.get("target", "")
                kpi_value_p.font.size = Pt(12)
                kpi_value_p.font.bold = True
                kpi_value_p.font.color.rgb = color

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_campaign_slide(
        self,
        title: str,
        campaign_name: str,
        period: str,
        objective: str,
        activities: list,
        notes: str = "",
    ):
        """
        캠페인 슬라이드 - 캠페인 개요 및 활동
        Modern 스타일: 헤더 배너, 활동 타임라인
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 캠페인 헤더 배너
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.4),
            Inches(CONTENT_WIDTH),
            Inches(1.2),
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        header_bg.line.fill.background()

        # 캠페인명
        campaign_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(1.5),
            Inches(5.0),
            Inches(0.5),
        )
        campaign_tf = campaign_box.text_frame
        campaign_p = campaign_tf.paragraphs[0]
        campaign_p.text = campaign_name
        campaign_p.font.size = Pt(22)
        campaign_p.font.bold = True
        campaign_p.font.color.rgb = RGBColor(255, 255, 255)

        # 기간
        period_box = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(2.0),
            Inches(4.0),
            Inches(0.4),
        )
        period_tf = period_box.text_frame
        period_p = period_tf.paragraphs[0]
        period_p.text = f"📅 {period}"
        period_p.font.size = Pt(12)
        period_p.font.color.rgb = STYLE_COLORS["sky_blue"]

        # 목표
        obj_box = slide.shapes.add_textbox(
            Inches(6.0),
            Inches(1.55),
            Inches(3.60),
            Inches(0.9),
        )
        obj_tf = obj_box.text_frame
        obj_tf.word_wrap = True
        obj_p = obj_tf.paragraphs[0]
        obj_p.text = objective
        obj_p.font.size = Pt(12)
        obj_p.font.color.rgb = RGBColor(200, 200, 200)
        obj_p.alignment = PP_ALIGN.RIGHT

        # 활동 타임라인/리스트
        activity_start_y = 2.9
        for i, activity in enumerate(activities[:6]):
            # 번호 원
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.7),
                Inches(activity_start_y + i * 0.75),
                Inches(0.4),
                Inches(0.4),
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = STYLE_COLORS["sky_blue"]
            num_circle.line.fill.background()

            # 번호 텍스트
            num_box = slide.shapes.add_textbox(
                Inches(0.7),
                Inches(activity_start_y + i * 0.75 + 0.05),
                Inches(0.4),
                Inches(0.35),
            )
            num_tf = num_box.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = str(i + 1)
            num_p.font.size = Pt(12)
            num_p.font.bold = True
            num_p.font.color.rgb = RGBColor(255, 255, 255)
            num_p.alignment = PP_ALIGN.CENTER

            # 활동 내용
            act_box = slide.shapes.add_textbox(
                Inches(1.3),
                Inches(activity_start_y + i * 0.75 + 0.05),
                Inches(8.25),
                Inches(0.5),
            )
            act_tf = act_box.text_frame
            act_tf.word_wrap = True
            act_p = act_tf.paragraphs[0]

            if isinstance(activity, dict):
                act_p.text = activity.get("name", str(activity))
            else:
                act_p.text = str(activity)

            act_p.font.size = Pt(14)
            act_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_budget_slide(
        self,
        title: str,
        budget_items: list,
        total: str = "",
        notes: str = "",
    ):
        """
        예산 슬라이드 - 투자 비용 테이블
        Modern 스타일: 깔끔한 테이블, 총계 강조
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 테이블 생성
        rows = len(budget_items) + 2  # 헤더 + 항목들 + 총계
        cols = 4  # 항목, 단가, 수량, 금액

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.8),
            Inches(1.6),
            Inches(8.80),
            Inches(0.5 * rows),
        )
        table = table_shape.table

        # 헤더 스타일
        headers = ["항목", "단가", "수량", "금액"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER

        # 데이터 행
        for row_idx, item in enumerate(budget_items, start=1):
            # 항목명
            cell0 = table.cell(row_idx, 0)
            cell0.text = item.get("name", "")
            cell0.text_frame.paragraphs[0].font.size = Pt(11)

            # 단가
            cell1 = table.cell(row_idx, 1)
            cell1.text = item.get("unit_price", "")
            cell1.text_frame.paragraphs[0].font.size = Pt(11)
            cell1.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # 수량
            cell2 = table.cell(row_idx, 2)
            cell2.text = str(item.get("quantity", ""))
            cell2.text_frame.paragraphs[0].font.size = Pt(11)
            cell2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 금액
            cell3 = table.cell(row_idx, 3)
            cell3.text = item.get("amount", "")
            cell3.text_frame.paragraphs[0].font.size = Pt(11)
            cell3.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 총계 행
        total_row = rows - 1
        table.cell(total_row, 0).merge(table.cell(total_row, 2))
        total_cell = table.cell(total_row, 0)
        total_cell.text = "총계"
        total_cell.fill.solid()
        total_cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        total_cell.text_frame.paragraphs[0].font.size = Pt(12)
        total_cell.text_frame.paragraphs[0].font.bold = True
        total_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        total_amount_cell = table.cell(total_row, 3)
        total_amount_cell.text = total
        total_amount_cell.fill.solid()
        total_amount_cell.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        total_amount_cell.text_frame.paragraphs[0].font.size = Pt(14)
        total_amount_cell.text_frame.paragraphs[0].font.bold = True
        total_amount_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        total_amount_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def add_case_study_slide(
        self,
        title: str,
        case: dict,
        notes: str = "",
    ):
        """
        케이스 스터디 슬라이드 - 수행 실적
        Modern 스타일: 프로젝트 이미지 + 성과 KPI
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 메인 제목
        self._add_title_textbox(slide, title)

        # 왼쪽: 이미지 플레이스홀더 + 프로젝트 정보
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5),
            Inches(1.5),
            Inches(4.88),
            Inches(4.0),
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        img_placeholder.line.fill.background()

        # 이미지 아이콘
        icon_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(3.2),
            Inches(4.88),
            Inches(0.5),
        )
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = "📷 프로젝트 이미지"
        icon_p.font.size = Pt(16)
        icon_p.font.color.rgb = RGBColor(150, 150, 150)
        icon_p.alignment = PP_ALIGN.CENTER

        # 프로젝트명
        proj_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(5.6),
            Inches(4.88),
            Inches(0.5),
        )
        proj_tf = proj_box.text_frame
        proj_p = proj_tf.paragraphs[0]
        proj_p.text = case.get("project_name", "")
        proj_p.font.size = Pt(18)
        proj_p.font.bold = True
        proj_p.font.color.rgb = STYLE_COLORS["dark_blue"]

        # 클라이언트/기간
        info_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(4.57),
            Inches(4.88),
            Inches(0.8),
        )
        info_tf = info_box.text_frame
        info_p = info_tf.paragraphs[0]
        client = case.get("client", "")
        period = case.get("period", "")
        info_p.text = f"{client} | {period}"
        info_p.font.size = Pt(12)
        info_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 오른쪽: 성과 KPI 영역
        kpi_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.4),
            Inches(1.5),
            Inches(4.22),
            Inches(5.5),
        )
        kpi_bg.fill.solid()
        kpi_bg.fill.fore_color.rgb = STYLE_COLORS["dark_blue"]
        kpi_bg.line.fill.background()

        # 성과 제목
        result_title_box = slide.shapes.add_textbox(
            Inches(5.55),
            Inches(1.7),
            Inches(5.23),
            Inches(0.5),
        )
        result_title_tf = result_title_box.text_frame
        result_title_p = result_title_tf.paragraphs[0]
        result_title_p.text = "주요 성과"
        result_title_p.font.size = Pt(16)
        result_title_p.font.bold = True
        result_title_p.font.color.rgb = RGBColor(255, 255, 255)

        # KPI 항목들
        kpis = case.get("kpis", case.get("results", []))
        for i, kpi in enumerate(kpis[:4]):
            kpi_box = slide.shapes.add_textbox(
                Inches(5.55),
                Inches(2.4 + i * 1.2),
                Inches(5.23),
                Inches(1.0),
            )
            kpi_tf = kpi_box.text_frame

            # KPI 값 (큰 숫자)
            kpi_value_p = kpi_tf.paragraphs[0]
            if isinstance(kpi, dict):
                kpi_value_p.text = kpi.get("value", kpi.get("target", ""))
            else:
                kpi_value_p.text = str(kpi)
            kpi_value_p.font.size = Pt(28)
            kpi_value_p.font.bold = True
            kpi_value_p.font.color.rgb = STYLE_COLORS["sky_blue"]

            # KPI 이름
            kpi_name_p = kpi_tf.add_paragraph()
            if isinstance(kpi, dict):
                kpi_name_p.text = kpi.get("name", kpi.get("metric", ""))
            else:
                kpi_name_p.text = ""
            kpi_name_p.font.size = Pt(11)
            kpi_name_p.font.color.rgb = RGBColor(200, 200, 200)

        # 프로젝트 설명 (아래)
        desc = case.get("description", case.get("overview", ""))
        if desc:
            desc_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(5.03),
                Inches(CONTENT_WIDTH),
                Inches(0.6),
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_p = desc_tf.paragraphs[0]
            desc_p.text = desc[:150] + "..." if len(desc) > 150 else desc
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = STYLE_COLORS["text_gray"]

        # 발표자 노트 추가
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # ========================================
    # v3.1 추가: Executive Summary, Next Step, Differentiation
    # ========================================

    def add_executive_summary_slide(
        self,
        project_objective: str,
        win_themes: list,
        kpis: list,
        why_us_points: list,
        notes: str = "",
    ):
        """
        Executive Summary 슬라이드 - 의사결정권자용 1페이지 요약

        Args:
            project_objective: 프로젝트 핵심 목표 (1문장)
            win_themes: [{"name": "Win Theme명", "description": "설명"}, ...]
            kpis: [{"metric": "KPI명", "target": "목표값", "basis": "산출근거"}, ...]
            why_us_points: ["포인트1", "포인트2", ...]
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 왼쪽 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(SLIDE_HEIGHT_INCHES)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = STYLE_COLORS["primary"]
        accent_bar.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6.0), Inches(0.6))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = "EXECUTIVE SUMMARY"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["primary"]

        # 프로젝트 목표
        obj_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9.23), Inches(0.6)
        )
        obj_bg.fill.solid()
        obj_bg.fill.fore_color.rgb = STYLE_COLORS["primary"]
        obj_bg.line.fill.background()

        obj_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.93), Inches(0.4))
        obj_p = obj_box.text_frame.paragraphs[0]
        obj_p.text = project_objective
        obj_p.font.size = Pt(16)
        obj_p.font.bold = True
        obj_p.font.color.rgb = STYLE_COLORS["white"]

        # Win Themes (3개 카드)
        win_colors = [STYLE_COLORS["primary"], STYLE_COLORS["secondary"], STYLE_COLORS["teal"]]
        for i, theme in enumerate(win_themes[:3]):
            x = Inches(0.5 + i * 4.2)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(4.0), Inches(1.4)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = win_colors[i % 3]
            card.line.fill.background()

            name_box = slide.shapes.add_textbox(x, Inches(2.1), Inches(4.0), Inches(0.5))
            name_p = name_box.text_frame.paragraphs[0]
            name_p.text = theme.get("name", "")
            name_p.font.size = Pt(14)
            name_p.font.bold = True
            name_p.font.color.rgb = STYLE_COLORS["white"]
            name_p.alignment = PP_ALIGN.CENTER

            desc_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.6), Inches(3.8), Inches(0.7))
            desc_box.text_frame.word_wrap = True
            desc_p = desc_box.text_frame.paragraphs[0]
            desc_p.text = theme.get("description", "")
            desc_p.font.size = Pt(11)
            desc_p.font.color.rgb = STYLE_COLORS["white"]
            desc_p.alignment = PP_ALIGN.CENTER

        # KPI 카드 (4개)
        for i, kpi in enumerate(kpis[:4]):
            x = Inches(0.5 + i * 3.2)
            kpi_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.6), Inches(3.0), Inches(1.5)
            )
            kpi_card.fill.solid()
            kpi_card.fill.fore_color.rgb = STYLE_COLORS["light"]
            kpi_card.line.fill.background()

            metric_box = slide.shapes.add_textbox(x, Inches(3.7), Inches(3.0), Inches(0.35))
            metric_p = metric_box.text_frame.paragraphs[0]
            metric_p.text = kpi.get("metric", "")
            metric_p.font.size = Pt(14)
            metric_p.font.bold = True
            metric_p.font.color.rgb = STYLE_COLORS["primary"]
            metric_p.alignment = PP_ALIGN.CENTER

            target_box = slide.shapes.add_textbox(x, Inches(4.05), Inches(3.0), Inches(0.4))
            target_p = target_box.text_frame.paragraphs[0]
            target_p.text = kpi.get("target", "")
            target_p.font.size = Pt(18)
            target_p.font.bold = True
            target_p.font.color.rgb = STYLE_COLORS["text_dark"]
            target_p.alignment = PP_ALIGN.CENTER

            basis_box = slide.shapes.add_textbox(x, Inches(4.5), Inches(3.0), Inches(0.55))
            basis_box.text_frame.word_wrap = True
            basis_p = basis_box.text_frame.paragraphs[0]
            basis_p.text = kpi.get("basis", kpi.get("calculation_basis", ""))
            basis_p.font.size = Pt(9)
            basis_p.font.color.rgb = STYLE_COLORS["text_gray"]
            basis_p.alignment = PP_ALIGN.CENTER

        # Why Us
        why_text = "  ".join([f"✓ {p}" for p in why_us_points[:4]])
        why_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9.23), Inches(0.5))
        why_box.text_frame.word_wrap = True
        why_p = why_box.text_frame.paragraphs[0]
        why_p.text = why_text
        why_p.font.size = Pt(12)
        why_p.font.bold = True
        why_p.font.color.rgb = STYLE_COLORS["secondary"]
        why_p.alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_next_step_slide(
        self,
        headline: str,
        steps: list,
        call_to_action: list,
        contact_info: dict = None,
        notes: str = "",
    ):
        """
        Next Step 슬라이드 - 다음 단계 안내 / Call to Action

        Args:
            headline: 헤드라인
            steps: [{"title": "계약 체결", "date": "2026.03", "description": "계약 협의"}, ...]
            call_to_action: ["10개월간 브랜드 인지도 +20%p 달성", ...]
            contact_info: {"name": "담당자명", "phone": "전화번호", "email": "이메일"}
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 왼쪽 액센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(SLIDE_HEIGHT_INCHES)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = STYLE_COLORS["primary"]
        accent_bar.line.fill.background()

        # 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6.0), Inches(0.6))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = "NEXT STEP"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["primary"]

        # 헤드라인
        headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.5))
        headline_p = headline_box.text_frame.paragraphs[0]
        headline_p.text = headline
        headline_p.font.size = Pt(24)
        headline_p.font.bold = True
        headline_p.font.color.rgb = STYLE_COLORS["text_dark"]

        # Step 카드들
        step_colors = [STYLE_COLORS["primary"]] + [STYLE_COLORS["secondary"]] * 10
        for i, step in enumerate(steps[:4]):
            x = Inches(0.5 + i * 3.2)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.0), Inches(1.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = step_colors[min(i, len(step_colors) - 1)]
            card.line.fill.background()

            num_box = slide.shapes.add_textbox(x, Inches(1.9), Inches(3.0), Inches(0.3))
            num_p = num_box.text_frame.paragraphs[0]
            num_p.text = f"STEP {i + 1}"
            num_p.font.size = Pt(11)
            num_p.font.bold = True
            num_p.font.color.rgb = STYLE_COLORS["white"]
            num_p.alignment = PP_ALIGN.CENTER

            title_b = slide.shapes.add_textbox(x, Inches(2.2), Inches(3.0), Inches(0.4))
            title_p2 = title_b.text_frame.paragraphs[0]
            title_p2.text = step.get("title", "")
            title_p2.font.size = Pt(18)
            title_p2.font.bold = True
            title_p2.font.color.rgb = STYLE_COLORS["white"]
            title_p2.alignment = PP_ALIGN.CENTER

            date_box = slide.shapes.add_textbox(x, Inches(2.6), Inches(3.0), Inches(0.3))
            date_p = date_box.text_frame.paragraphs[0]
            date_p.text = step.get("date", "")
            date_p.font.size = Pt(12)
            date_p.font.color.rgb = STYLE_COLORS["white"]
            date_p.alignment = PP_ALIGN.CENTER

            desc_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.95), Inches(2.8), Inches(0.55))
            desc_box.text_frame.word_wrap = True
            desc_p = desc_box.text_frame.paragraphs[0]
            desc_p.text = step.get("description", "")
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = STYLE_COLORS["white"]
            desc_p.alignment = PP_ALIGN.CENTER

            if i < len(steps) - 1:
                arrow_box = slide.shapes.add_textbox(x + Inches(3.05), Inches(2.5), Inches(0.3), Inches(0.4))
                arrow_p = arrow_box.text_frame.paragraphs[0]
                arrow_p.text = "→"
                arrow_p.font.size = Pt(20)
                arrow_p.font.bold = True
                arrow_p.font.color.rgb = STYLE_COLORS["text_light"]

        # CTA 영역
        cta_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.9), Inches(9.23), Inches(1.2)
        )
        cta_bg.fill.solid()
        cta_bg.fill.fore_color.rgb = STYLE_COLORS["light"]
        cta_bg.line.fill.background()

        cta_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.93), Inches(0.35))
        cta_title_p = cta_title.text_frame.paragraphs[0]
        cta_title_p.text = "저희가 제안하는 것"
        cta_title_p.font.size = Pt(14)
        cta_title_p.font.bold = True
        cta_title_p.font.color.rgb = STYLE_COLORS["primary"]

        for i, cta in enumerate(call_to_action[:4]):
            x = Inches(0.9) if i < 2 else Inches(4.88)
            y = Inches(4.4 + (i % 2) * 0.35)
            cta_box = slide.shapes.add_textbox(x, y, Inches(5.4), Inches(0.35))
            cta_p = cta_box.text_frame.paragraphs[0]
            cta_p.text = f"✓ {cta}"
            cta_p.font.size = Pt(12)
            cta_p.font.color.rgb = STYLE_COLORS["text_dark"]

        # 연락처
        if contact_info:
            contact_text = f"담당자: {contact_info.get('name', '[담당자명]')} | 연락처: {contact_info.get('phone', '[전화번호]')} | 이메일: {contact_info.get('email', '[이메일]')}"
        else:
            contact_text = "담당자: [담당자명] | 연락처: [전화번호] | 이메일: [이메일]"

        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9.23), Inches(0.4))
        contact_p = contact_box.text_frame.paragraphs[0]
        contact_p.text = contact_text
        contact_p.font.size = Pt(12)
        contact_p.font.color.rgb = STYLE_COLORS["text_gray"]
        contact_p.alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_section_divider_with_win_theme(
        self,
        phase_number: int,
        phase_title: str,
        phase_subtitle: str = "",
        story_title: str = "",
        win_theme: str = "",
        notes: str = "",
    ):
        """
        섹션 구분 슬라이드 (Win Theme 배지 포함)
        """
        slide_layout = self._get_layout("Title and Content_02")
        slide = self._add_slide_with_placeholders(slide_layout)

        # 다크 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH_INCHES), Inches(SLIDE_HEIGHT_INCHES)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = STYLE_COLORS["dark_bg"]
        bg.line.fill.background()

        # Phase 번호
        num_text = f"0{phase_number}" if phase_number < 10 else str(phase_number)
        num_box = slide.shapes.add_textbox(Inches(5.0), Inches(0.5), Inches(6.0), Inches(4.5))
        num_p = num_box.text_frame.paragraphs[0]
        num_p.text = num_text
        num_p.font.size = Pt(320)
        num_p.font.bold = True
        num_p.font.color.rgb = RGBColor(40, 40, 40)
        num_p.alignment = PP_ALIGN.RIGHT

        # Part 라벨
        part_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(4), Inches(0.5))
        part_p = part_box.text_frame.paragraphs[0]
        part_p.text = f"PART {num_text}"
        part_p.font.size = Pt(16)
        part_p.font.bold = True
        part_p.font.color.rgb = STYLE_COLORS["secondary"]

        # 스토리 타이틀
        if story_title:
            story_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(6.0), Inches(0.5))
            story_p = story_box.text_frame.paragraphs[0]
            story_p.text = story_title
            story_p.font.size = Pt(20)
            story_p.font.color.rgb = STYLE_COLORS["secondary"]

        # 메인 타이틀
        y_title = Inches(3.9) if story_title else Inches(3.6)
        title_box = slide.shapes.add_textbox(Inches(0.8), y_title, Inches(6.0), Inches(1.0))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = phase_title
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = STYLE_COLORS["white"]

        # 서브타이틀
        if phase_subtitle:
            y_sub = Inches(5.3) if story_title else Inches(5.0)
            sub_box = slide.shapes.add_textbox(Inches(0.8), y_sub, Inches(6.0), Inches(0.5))
            sub_p = sub_box.text_frame.paragraphs[0]
            sub_p.text = phase_subtitle
            sub_p.font.size = Pt(16)
            sub_p.font.color.rgb = STYLE_COLORS["text_light"]

        # Win Theme 배지
        if win_theme:
            badge_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.88), Inches(6.0), Inches(0.6)
            )
            badge_bg.fill.solid()
            badge_bg.fill.fore_color.rgb = STYLE_COLORS["secondary"]
            badge_bg.line.fill.background()

            badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.94), Inches(6.0), Inches(0.45))
            badge_p = badge_box.text_frame.paragraphs[0]
            badge_p.text = f"💡 Win Theme: {win_theme}"
            badge_p.font.size = Pt(14)
            badge_p.font.bold = True
            badge_p.font.color.rgb = STYLE_COLORS["white"]

        if notes:
            slide.notes_slide.notes_text_frame.text = notes
