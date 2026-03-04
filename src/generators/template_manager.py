"""
PPTX 템플릿 관리자 ([회사명])

템플릿 PPTX 로드, 레이아웃 인덱스 관리, 디자인 시스템(색상·폰트·간격) 제공.
로드한 템플릿 파일의 테마(색상·폰트)와 슬라이드 크기·플레이스홀더 위치를 동적으로 추출해
제안서 작성 시 하드코딩 없이 적용합니다. template_name이 비어 있으면 템플릿 미사용(빈 프레젠테이션 + 기본 디자인).
"""

import json
from pathlib import Path
from typing import Any, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..utils.logger import get_logger
from ..utils.path_utils import safe_filename
from config.settings import get_settings

logger = get_logger("template_manager")
_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class TemplateManager:
    """
    PPTX 템플릿 관리자.

    templates_dir 내 .pptx 로드, slide_layouts.json 기반 레이아웃 인덱스 제공,
    테마 추출로 design_system(colors, fonts, spacing) 갱신, 플레이스홀더 기하(_layout_geometry) 추출.
    """

    def __init__(self, templates_dir: Optional[Path] = None):
        settings = get_settings()
        self.templates_dir = templates_dir or settings.templates_dir
        self.layouts = self._load_layouts()
        self.design_system = self._get_design_system()
        self._cached_guide_path: Optional[Path] = None
        # 템플릿 PPTX에서 동적 추출한 레이아웃/위치 (하드코딩 대체)
        self._layout_geometry: Optional[Dict[str, Any]] = None

    def _load_layouts(self) -> Dict[str, Any]:
        """레이아웃 정의 로드"""
        layout_file = self.templates_dir / "slide_layouts.json"
        if layout_file.exists():
            try:
                return json.loads(layout_file.read_text(encoding="utf-8"))
            except Exception as e:
                logger.warning(f"레이아웃 파일 로드 실패: {e}")

        return self._default_layouts()

    def _default_layouts(self) -> Dict[str, Any]:
        """기본 레이아웃 정의"""
        return {
            "layouts": {
                "title": {"index": 0, "name": "Title Slide"},
                "section": {"index": 2, "name": "Section Header"},
                "content": {"index": 1, "name": "Title and Content"},
                "two_column": {"index": 3, "name": "Two Content"},
                "comparison": {"index": 4, "name": "Comparison"},
                "blank": {"index": 6, "name": "Blank"},
            }
        }

    def _hex_to_rgb(self, value: Any) -> RGBColor:
        """#RRGGBB 또는 [r,g,b] → RGBColor"""
        if isinstance(value, list) and len(value) >= 3:
            return RGBColor(value[0], value[1], value[2])
        s = str(value).strip()
        if s.startswith("#"):
            s = s[1:]
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        return RGBColor(51, 51, 51)

    def _get_default_design_system(self) -> Dict[str, Any]:
        """기본 디자인 시스템. 템플릿 미로드 시 또는 PPTX 테마 추출 실패 시 사용."""
        return {
            "colors": {
                "primary": RGBColor(0, 82, 147),
                "secondary": RGBColor(0, 150, 199),
                "accent": RGBColor(255, 107, 0),
                "success": RGBColor(40, 167, 69),
                "warning": RGBColor(255, 193, 7),
                "danger": RGBColor(220, 53, 69),
                "text_dark": RGBColor(51, 51, 51),
                "text_light": RGBColor(128, 128, 128),
                "text_gray": RGBColor(102, 102, 102),
                "background": RGBColor(255, 255, 255),
                "background_light": RGBColor(245, 245, 245),
                "white": RGBColor(255, 255, 255),
                "light": RGBColor(245, 245, 245),
                "dark_blue": RGBColor(0, 44, 95),
                "sky_blue": RGBColor(0, 170, 210),
                "dark_bg": RGBColor(26, 26, 26),
                "teal": RGBColor(0, 161, 156),
            },
            "fonts": {
                "title": "맑은 고딕",
                "body": "맑은 고딕",
                "english": "Arial",
                "sizes": {
                    "cover_title": Pt(44),
                    "part_title": Pt(40),
                    "slide_title": Pt(28),
                    "subtitle": Pt(20),
                    "body": Pt(16),
                    "small": Pt(14),
                    "caption": Pt(12),
                },
            },
            "spacing": {
                "margin": Inches(0.5),
                "content_margin": Inches(0.75),
                "element_gap": Inches(0.25),
            },
            "dimensions": {
                "slide_width": Inches(13.33),
                "slide_height": Inches(7.5),
            },
        }

    def _get_design_system(self) -> Dict[str, Any]:
        """초기 디자인 시스템. load_template()에서 템플릿 PPTX 기준으로 갱신됨."""
        return self._get_default_design_system()

    def _extract_design_from_presentation(self, prs: Presentation) -> Optional[Dict[str, Any]]:
        """
        로드된 PPTX의 테마(색상·폰트)를 추출. 반환값은 기본 디자인과 병합해 사용.
        실패 시 None.
        """
        try:
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            from pptx.oxml import parse_xml
        except ImportError:
            return None
        try:
            theme_part = prs.part.part_related_by(RT.THEME)
        except Exception:
            return None
        if not theme_part or not getattr(theme_part, "blob", None):
            return None
        root = parse_xml(theme_part.blob)
        if root is None:
            return None
        out: Dict[str, Any] = {"colors": {}, "fonts": {}}
        scheme = root.find(f".//{{{_DML_NS}}}clrScheme")
        if scheme is not None:
            for child in scheme:
                tag = child.tag.replace(f"{{{_DML_NS}}}", "") if "}" in child.tag else child.tag
                srgb = child.find(f"{{{_DML_NS}}}srgbClr")
                if srgb is None or not srgb.get("val"):
                    continue
                hex_val = srgb.get("val", "").strip()
                rgb = self._hex_to_rgb(hex_val)
                if tag == "dk1":
                    out["colors"]["text_dark"] = rgb
                elif tag == "lt1":
                    out["colors"]["background"] = out["colors"]["white"] = rgb
                elif tag == "dk2":
                    out["colors"]["text_gray"] = rgb
                elif tag == "lt2":
                    out["colors"]["background_light"] = out["colors"]["light"] = rgb
                elif tag == "accent1":
                    out["colors"]["primary"] = out["colors"]["dark_blue"] = rgb
                elif tag == "accent2":
                    out["colors"]["secondary"] = out["colors"]["sky_blue"] = rgb
                elif tag == "accent3":
                    out["colors"]["teal"] = rgb
                elif tag == "accent4":
                    out["colors"]["danger"] = rgb
                elif tag == "accent5":
                    out["colors"]["dark_bg"] = rgb
                elif tag == "accent6":
                    out["colors"]["accent"] = rgb
        font_scheme = root.find(f".//{{{_DML_NS}}}fontScheme")
        if font_scheme is not None:
            for role, key in (("majorFont", "title"), ("minorFont", "body")):
                node = font_scheme.find(f"{{{_DML_NS}}}{role}")
                if node is None:
                    continue
                latin = node.find(f"{{{_DML_NS}}}latin")
                if latin is not None and latin.get("typeface"):
                    name = latin.get("typeface", "").strip()
                    out["fonts"][key] = name
            if out["fonts"].get("title") and not out["fonts"].get("body"):
                out["fonts"]["body"] = out["fonts"]["title"]
        if out["colors"] or out["fonts"]:
            return out
        return None

    def _apply_design_from_presentation(self, prs: Presentation) -> None:
        """
        로드된 PPTX에서 테마(색상·폰트)를 추출해 self.design_system에 반영.
        추출 실패 시 기본값 유지.
        """
        extracted = self._extract_design_from_presentation(prs)
        if not extracted:
            return
        if extracted.get("colors"):
            for k, v in extracted["colors"].items():
                self.design_system["colors"][k] = v
            logger.info("템플릿 PPTX 테마 색상 적용: {}", list(extracted["colors"].keys()))
        if extracted.get("fonts"):
            for k in ("title", "body", "english"):
                if k in extracted["fonts"] and extracted["fonts"][k]:
                    self.design_system["fonts"][k] = str(extracted["fonts"][k])
            logger.info("템플릿 PPTX 테마 폰트 적용: {}", extracted.get("fonts"))

    def _extract_layout_from_presentation(self, prs: Presentation) -> None:
        """
        템플릿 PPTX의 슬라이드 크기·플레이스홀더 위치를 추출해 _layout_geometry에 저장.
        생성 시 하드코딩 대신 이 값만 사용하도록 함.
        """
        try:
            def to_inch(v) -> float:
                if v is None:
                    return 0.0
                return getattr(v, "inches", None) or (float(v) / 914400.0)
            sw = to_inch(prs.slide_width)
            sh = to_inch(prs.slide_height)
            placeholders: Dict[str, Dict[str, Any]] = {}
            for layout in prs.slide_layouts:
                for ph in layout.placeholders:
                    idx = ph.placeholder_format.idx
                    role = "title" if idx == 0 else ("body" if idx == 1 else f"idx_{idx}")
                    if role in placeholders:
                        continue
                    geom = {
                        "left": to_inch(ph.left),
                        "top": to_inch(ph.top),
                        "width": to_inch(ph.width),
                        "height": to_inch(ph.height),
                    }
                    if hasattr(ph, "text_frame") and ph.text_frame.paragraphs:
                        p0 = ph.text_frame.paragraphs[0]
                        if getattr(p0.font, "size", None) is not None:
                            geom["font_size_pt"] = p0.font.size.pt
                        if getattr(p0.font, "name", None):
                            geom["font_name"] = p0.font.name
                    placeholders[role] = geom
                if placeholders:
                    break
            self._layout_geometry = {
                "slide_width_inches": sw,
                "slide_height_inches": sh,
                "placeholders": placeholders,
            }
            logger.info(
                "템플릿 레이아웃 추출: 슬라이드 {:.2f} x {:.2f} inch, 플레이스홀더 {}",
                sw, sh, list(placeholders.keys()),
            )
        except Exception as e:
            logger.warning("템플릿 레이아웃 추출 실패: {}", e)
            self._layout_geometry = None

    def get_slide_width_inches(self) -> Optional[float]:
        """템플릿에서 추출한 슬라이드 너비(inch). 없으면 None."""
        if not self._layout_geometry:
            return None
        return self._layout_geometry.get("slide_width_inches")

    def get_slide_height_inches(self) -> Optional[float]:
        """템플릿에서 추출한 슬라이드 높이(inch). 없으면 None."""
        if not self._layout_geometry:
            return None
        return self._layout_geometry.get("slide_height_inches")

    def get_placeholder_geometry(self, role: str) -> Optional[Dict[str, Any]]:
        """
        템플릿에서 추출한 플레이스홀더 위치/크기. role: 'title', 'body' 등.
        반환: { left, top, width, height (, font_size_pt, font_name) } 또는 None.
        """
        if not self._layout_geometry:
            return None
        return (self._layout_geometry.get("placeholders") or {}).get(role)

    def _find_guide_template(self) -> Optional[Path]:
        """
        templates 폴더 하위에서 '가이드' 또는 'guide'가 포함된 .pptx 파일을 찾음.
        결과를 인스턴스에 캐시하여 반복 glob 방지.
        """
        if self._cached_guide_path is not None:
            return self._cached_guide_path
        if not self.templates_dir.exists():
            return None
        keywords = ("가이드", "guide")
        for f in sorted(self.templates_dir.glob("*.pptx")):
            stem_lower = f.stem.lower()
            if any(kw in stem_lower or kw in f.stem for kw in keywords):
                self._cached_guide_path = f
                return f
        return None

    def load_template(self, template_name: Optional[str] = "base_template") -> Presentation:
        """
        템플릿 PPTX 로드 후 해당 파일의 테마(색상·폰트)를 제안서 작성 규칙으로 적용.

        template_name이 None 또는 빈 문자열이면: 템플릿 파일을 사용하지 않고
        빈 프레젠테이션 + 기본 디자인 시스템만 사용 (기존 품질 우선 로직).

        그 외:
        1) templates/{template_name}.pptx 가 있으면 해당 파일 사용
        2) 없으면 templates/ 하위에서 '가이드' 또는 'guide'가 포함된 .pptx를 동적으로 검색해 사용
        3) 둘 다 없으면 빈 프레젠테이션 생성 (기본 디자인 시스템 사용)

        Args:
            template_name: None/빈 문자열 = 템플릿 미사용(기본 디자인). 파일명 지정 시 해당 .pptx 사용.

        Returns:
            Presentation 객체 (슬라이드 0장, 레이아웃/테마만 적용된 상태)
        """
        if not (template_name or "").strip():
            logger.info("템플릿 미사용: 기본 디자인으로 제안서 생성 (빈 프레젠테이션)")
            self._layout_geometry = None
            return Presentation()

        safe_name = safe_filename(template_name, max_len=80)
        template_path = (self.templates_dir / f"{safe_name}.pptx").resolve()
        templates_resolved = self.templates_dir.resolve()
        try:
            template_path.relative_to(templates_resolved)
        except ValueError:
            template_path = templates_resolved / f"{safe_name}.pptx"

        if template_path.exists():
            logger.info("템플릿 로드: {} (해당 PPTX 테마·폰트·레이아웃 기준으로 제안서 생성)", template_path.name)
            prs = Presentation(template_path)
            self._clear_all_slides(prs)
            self._apply_design_from_presentation(prs)
            self._extract_layout_from_presentation(prs)
            return prs

        guide_path = self._find_guide_template()
        if guide_path is not None:
            logger.info("가이드 템플릿 로드: {} (해당 PPTX 테마·폰트·레이아웃 기준으로 제안서 생성)", guide_path.name)
            prs = Presentation(guide_path)
            self._clear_all_slides(prs)
            self._apply_design_from_presentation(prs)
            self._extract_layout_from_presentation(prs)
            return prs

        logger.info("기본 빈 프레젠테이션 생성 (templates 내 guide 포함 .pptx 없음)")
        self._layout_geometry = None
        return Presentation()

    def _clear_all_slides(self, prs: Presentation) -> None:
        """
        템플릿의 기존 슬라이드 전부 제거. 레이아웃·테마는 유지하여
        목차 형식 등 디자인만 참조하고 생성 콘텐츠만 넣을 때 사용.
        """
        try:
            sld_id_lst = prs.slides._sldIdLst
            for i in range(len(prs.slides) - 1, -1, -1):
                r_id = sld_id_lst[i].rId
                prs.part.drop_rel(r_id)
                del sld_id_lst[i]
            logger.info("템플릿 기존 슬라이드 제거 완료 (레이아웃/테마만 사용)")
        except Exception as e:
            logger.warning("템플릿 슬라이드 제거 중 오류: {}", e)

    def get_layout_index(self, layout_name: str) -> int:
        """
        레이아웃 인덱스 반환

        Args:
            layout_name: 레이아웃 이름

        Returns:
            레이아웃 인덱스 (기본값: 1)
        """
        layouts = self.layouts.get("layouts", {})
        layout = layouts.get(layout_name, {})
        return layout.get("index", 1)

    def get_color(self, color_name: str) -> RGBColor:
        """색상 반환"""
        return self.design_system["colors"].get(
            color_name, self.design_system["colors"]["text_dark"]
        )

    def get_font_size(self, size_name: str) -> Pt:
        """폰트 크기 반환"""
        return self.design_system["fonts"]["sizes"].get(
            size_name, self.design_system["fonts"]["sizes"]["body"]
        )

    def get_font_size_for_text(self, text: str, base_size_name: str = "body") -> Pt:
        """고도화: 텍스트 길이에 따라 폰트 크기 자동 조절 (긴 텍스트는 작게)."""
        if not (text or "").strip():
            return self.get_font_size(base_size_name)
        n = len(str(text).strip())
        if n > 800:
            return self.get_font_size("caption")
        if n > 400:
            return self.get_font_size("small")
        return self.get_font_size(base_size_name)

    def get_font_name(self, font_type: str = "body") -> str:
        """폰트 이름 반환"""
        return self.design_system["fonts"].get(font_type, "맑은 고딕")
